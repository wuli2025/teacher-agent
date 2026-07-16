# -*- coding: utf-8 -*-
"""American-Keynote-style deck engine on top of python-pptx.

Design language: full-bleed hero imagery, dark scrims, oversized type,
generous whitespace, a strict per-deck color system, section dividers.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_IN = 914400
PW, PH = 13.333, 7.5   # 16:9 inches


def _rgb(hexs):
    return RGBColor.from_string(hexs.replace("#", ""))


class Theme:
    def __init__(self, ink, paper, accent, accent2, muted, fontEN, fontHead, fontBody):
        self.ink = _rgb(ink)
        self.paper = _rgb(paper)
        self.accent = _rgb(accent)
        self.accent2 = _rgb(accent2)
        self.muted = _rgb(muted)
        self.fontEN = fontEN
        self.fontHead = fontHead
        self.fontBody = fontBody


class Deck:
    def __init__(self, theme):
        self.prs = Presentation()
        self.prs.slide_width = Inches(PW)
        self.prs.slide_height = Inches(PH)
        self.t = theme
        self.blank = self.prs.slide_layouts[6]

    # ---- low level helpers ----
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _bg(self, s, color):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = color

    def _rect(self, s, x, y, w, h, color, alpha=None, line=False):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        if not line:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = color
        sp.shadow.inherit = False
        if alpha is not None:
            self._alpha(sp, alpha)
        return sp

    def _alpha(self, sp, alpha):
        # alpha 0..100 (transparency percent)
        el = sp.fill._xPr.find(qn('a:solidFill'))
        srgb = el.find(qn('a:srgbClr'))
        a = srgb.makeelement(qn('a:alpha'), {'val': str(int((100 - alpha) * 1000))})
        srgb.append(a)

    def _img(self, s, path, x, y, w, h):
        # cover-fit crop into box
        from PIL import Image
        iw, ih = Image.open(path).size
        boxr = w / h
        imgr = iw / ih
        pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        if imgr > boxr:
            crop = (1 - boxr / imgr) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - imgr / boxr) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
        return pic

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              space=1.0, wrap=True):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = para.get("align", align)
            p.line_spacing = para.get("space", space)
            if "before" in para:
                p.space_before = Pt(para["before"])
            if "after" in para:
                p.space_after = Pt(para["after"])
            for seg in para["segs"]:
                r = p.add_run()
                r.text = seg["t"]
                f = r.font
                f.size = Pt(seg.get("sz", 18))
                f.name = seg.get("font", self.t.fontBody)
                f.bold = seg.get("b", False)
                f.italic = seg.get("i", False)
                f.color.rgb = seg.get("c", self.t.ink)
                # set east-asian font too
                rPr = r._r.get_or_add_rPr()
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = rPr.makeelement(qn('a:ea'), {})
                    rPr.append(ea)
                ea.set('typeface', seg.get("font", self.t.fontBody))
                if seg.get("spacing"):
                    rPr.set('spc', str(int(seg["spacing"] * 100)))
        return tb

    def _line(self, s, x, y, w, color, weight=2.2):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(weight))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def _pageno(self, s, n, dark=False):
        c = self.t.paper if dark else self.t.muted
        self._text(s, PW - 1.2, PH - 0.55, 0.8, 0.4,
                   [{"segs": [{"t": f"{n:02d}", "sz": 11, "c": c, "font": self.t.fontEN,
                               "spacing": 2}]}], align=PP_ALIGN.RIGHT)

    # ---- high level layouts ----
    def title_hero(self, img, kicker, title_lines, subtitle, footer):
        s = self._slide()
        self._img(s, img, 0, 0, PW, PH)
        self._rect(s, 0, 0, PW, PH, self.t.ink, alpha=42)
        self._rect(s, 0, PH - 2.9, PW, 2.9, self.t.ink, alpha=22)
        y = 2.15
        self._line(s, 1.0, y, 0.9, self.t.accent, weight=3)
        self._text(s, 1.0, y + 0.12, 11, 0.5,
                   [{"segs": [{"t": kicker, "sz": 15, "c": self.t.paper, "b": True,
                               "font": self.t.fontEN, "spacing": 5}]}])
        segs = []
        runs = []
        for ln in title_lines:
            runs.append({"segs": [{"t": ln, "sz": 54, "c": self.t.paper, "b": True,
                                   "font": self.t.fontHead}], "space": 1.02})
        self._text(s, 0.97, y + 0.62, 11.5, 3.2, runs)
        self._text(s, 1.0, PH - 1.55, 10.5, 0.7,
                   [{"segs": [{"t": subtitle, "sz": 18, "c": self.t.paper,
                               "font": self.t.fontBody}]}])
        self._text(s, 1.0, PH - 0.72, 11.3, 0.4,
                   [{"segs": [{"t": footer, "sz": 12, "c": self.t.paper,
                               "font": self.t.fontEN, "spacing": 3}]}])
        return s

    def divider(self, img, num, title, subtitle):
        s = self._slide()
        self._img(s, img, 0, 0, PW, PH)
        self._rect(s, 0, 0, PW, PH, self.t.ink, alpha=48)
        self._rect(s, 0, 0, 0.35, PH, self.t.accent)
        self._text(s, 1.1, 2.35, 4, 1.6,
                   [{"segs": [{"t": num, "sz": 96, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN}]}])
        self._line(s, 1.15, 4.15, 1.6, self.t.paper, weight=2.5)
        self._text(s, 1.13, 4.35, 11, 1.4,
                   [{"segs": [{"t": title, "sz": 44, "c": self.t.paper, "b": True,
                               "font": self.t.fontHead}]}])
        if subtitle:
            self._text(s, 1.15, 5.55, 10.5, 0.8,
                       [{"segs": [{"t": subtitle, "sz": 17, "c": self.t.paper,
                                   "font": self.t.fontBody}]}])
        return s

    def image_text(self, img, kicker, title, body_paras, side="left", pageno=None):
        s = self._slide()
        self._bg(s, self.t.paper)
        halfw = 5.85
        if side == "left":
            self._img(s, img, 0, 0, halfw, PH)
            tx = halfw + 0.85
        else:
            self._img(s, img, PW - halfw, 0, halfw, PH)
            tx = 0.85
        tw = PW - halfw - 1.7
        self._text(s, tx, 1.15, tw, 0.4,
                   [{"segs": [{"t": kicker, "sz": 13, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN, "spacing": 4}]}])
        self._text(s, tx, 1.5, tw, 1.5,
                   [{"segs": [{"t": title, "sz": 33, "c": self.t.ink, "b": True,
                               "font": self.t.fontHead}], "space": 1.05}])
        self._line(s, tx + 0.02, 2.55, 0.7, self.t.accent, weight=3)
        runs = []
        for para in body_paras:
            runs.append({"segs": [{"t": para, "sz": 16.5, "c": self.t.ink,
                                   "font": self.t.fontBody}], "space": 1.28, "after": 10})
        self._text(s, tx, 2.9, tw, 4, runs)
        if pageno:
            self._pageno(s, pageno)
        return s

    def bullets(self, kicker, title, items, img=None, pageno=None):
        s = self._slide()
        self._bg(s, self.t.paper)
        self._rect(s, 0, 0, PW, 0.28, self.t.accent)
        tw = 7.0 if img else 11.3
        self._text(s, 1.0, 0.9, tw, 0.4,
                   [{"segs": [{"t": kicker, "sz": 13, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN, "spacing": 4}]}])
        self._text(s, 1.0, 1.25, tw, 1.1,
                   [{"segs": [{"t": title, "sz": 32, "c": self.t.ink, "b": True,
                               "font": self.t.fontHead}]}])
        self._line(s, 1.02, 2.25, 0.7, self.t.accent, weight=3)
        # single column (with image, or few items) vs. two columns (many items, no image)
        two_col = (img is None and len(items) > 4)
        step = 1.08  # vertical space per item with a sub-line
        if two_col:
            half = (len(items) + 1) // 2
            columns = [(1.0, items[:half]), (7.05, items[half:])]
            iw = 5.1
        else:
            columns = [(1.0, items)]
            iw = tw - 0.4

        for cx, col_items in columns:
            y = 2.75
            for it in col_items:
                head = it[0]; sub = it[1] if len(it) > 1 else None
                self._rect(s, cx, y + 0.13, 0.16, 0.16, self.t.accent2)
                runs = [{"segs": [{"t": head, "sz": 18, "c": self.t.ink, "b": True,
                                   "font": self.t.fontBody}]}]
                if sub:
                    runs.append({"segs": [{"t": sub, "sz": 14, "c": self.t.muted,
                                           "font": self.t.fontBody}], "space": 1.2, "before": 3})
                self._text(s, cx + 0.4, y, iw, 1.2, runs)
                y += 0.6 + (step - 0.6 if sub else 0)
        if img:
            self._img(s, img, 8.35, 1.15, 4.15, 5.2)
        if pageno:
            self._pageno(s, pageno)
        return s

    def fullbleed_caption(self, img, caption, sub=None, pageno=None):
        s = self._slide()
        self._img(s, img, 0, 0, PW, PH)
        self._rect(s, 0, PH - 2.2, PW, 2.2, self.t.ink, alpha=30)
        runs = [{"segs": [{"t": caption, "sz": 30, "c": self.t.paper, "b": True,
                           "font": self.t.fontHead}]}]
        if sub:
            runs.append({"segs": [{"t": sub, "sz": 16, "c": self.t.paper,
                                   "font": self.t.fontBody}], "before": 6})
        self._text(s, 1.0, PH - 1.85, 11.3, 1.5, runs, anchor=MSO_ANCHOR.TOP)
        self._line(s, 1.02, PH - 2.05, 0.7, self.t.accent, weight=3)
        if pageno:
            self._pageno(s, pageno, dark=True)
        return s

    def quote(self, text, cite, img=None):
        s = self._slide()
        if img:
            self._img(s, img, 0, 0, PW, PH)
            self._rect(s, 0, 0, PW, PH, self.t.ink, alpha=55)
            tc = self.t.paper
        else:
            self._bg(s, self.t.ink)
            tc = self.t.paper
        self._text(s, 1.6, 1.4, 10.1, 0.9,
                   [{"segs": [{"t": "“", "sz": 90, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN}]}])
        self._text(s, 1.6, 2.6, 10.1, 3,
                   [{"segs": [{"t": text, "sz": 30, "c": tc, "i": True,
                               "font": self.t.fontHead}], "space": 1.25, "align": PP_ALIGN.LEFT}])
        self._line(s, 1.62, 5.7, 0.7, self.t.accent, weight=3)
        self._text(s, 1.6, 5.9, 10, 0.6,
                   [{"segs": [{"t": cite, "sz": 15, "c": tc, "b": True,
                               "font": self.t.fontBody, "spacing": 2}]}])
        return s

    def cards(self, kicker, title, cards, cols=2, pageno=None):
        s = self._slide()
        self._bg(s, self.t.paper)
        self._rect(s, 0, 0, PW, 0.28, self.t.accent)
        self._text(s, 1.0, 0.75, 11, 0.4,
                   [{"segs": [{"t": kicker, "sz": 13, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN, "spacing": 4}]}])
        self._text(s, 1.0, 1.1, 11.3, 1.0,
                   [{"segs": [{"t": title, "sz": 30, "c": self.t.ink, "b": True,
                               "font": self.t.fontHead}]}])
        n = len(cards)
        rows = (n + cols - 1) // cols
        gx, gy = 0.4, 0.4
        x0, y0 = 1.0, 2.35
        cw = (PW - 2 * x0 - (cols - 1) * gx) / cols
        ch = (PH - y0 - 0.7 - (rows - 1) * gy) / rows
        for i, card in enumerate(cards):
            r, cc = divmod(i, cols)
            x = x0 + cc * (cw + gx)
            y = y0 + r * (ch + gy)
            # clean white card with a thin border and a colored top accent bar
            card_rect = self._rect(s, x, y, cw, ch, _rgb("FFFFFF"))
            card_rect.line.color.rgb = _rgb("E3DED3")
            card_rect.line.width = Pt(1)
            self._rect(s, x, y, cw, 0.12, self.t.accent)
            head = card[0]; body = card[1]
            self._text(s, x + 0.35, y + 0.34, cw - 0.6, 0.6,
                       [{"segs": [{"t": head, "sz": 20, "c": self.t.accent, "b": True,
                                   "font": self.t.fontBody}]}])
            self._text(s, x + 0.35, y + 0.98, cw - 0.6, ch - 1.15,
                       [{"segs": [{"t": body, "sz": 14.5, "c": self.t.ink,
                                   "font": self.t.fontBody}], "space": 1.24}])
        if pageno:
            self._pageno(s, pageno)
        return s

    def statement(self, big, sub=None, accent_bg=False):
        s = self._slide()
        self._bg(s, self.t.accent if accent_bg else self.t.ink)
        tc = self.t.paper
        self._text(s, 1.2, 0, 10.9, PH,
                   [{"segs": [{"t": big, "sz": 46, "c": tc, "b": True,
                               "font": self.t.fontHead}], "space": 1.1, "align": PP_ALIGN.LEFT}],
                   anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            self._text(s, 1.22, PH - 1.7, 10.5, 1,
                       [{"segs": [{"t": sub, "sz": 18, "c": tc, "font": self.t.fontBody}]}])
        return s

    def two_col(self, kicker, title, left_title, left_items, right_title, right_items, pageno=None):
        s = self._slide()
        self._bg(s, self.t.paper)
        self._rect(s, 0, 0, PW, 0.28, self.t.accent)
        self._text(s, 1.0, 0.75, 11, 0.4,
                   [{"segs": [{"t": kicker, "sz": 13, "c": self.t.accent, "b": True,
                               "font": self.t.fontEN, "spacing": 4}]}])
        self._text(s, 1.0, 1.1, 11.3, 1.0,
                   [{"segs": [{"t": title, "sz": 30, "c": self.t.ink, "b": True,
                               "font": self.t.fontHead}]}])
        cols = [(1.0, left_title, left_items, self.t.accent),
                (7.0, right_title, right_items, self.t.accent2)]
        for x, ct, items, col in cols:
            self._rect(s, x, 2.4, 5.3, 0.62, col)
            self._text(s, x + 0.3, 2.5, 5.0, 0.5,
                       [{"segs": [{"t": ct, "sz": 18, "c": self.t.paper, "b": True,
                                   "font": self.t.fontBody}]}])
            y = 3.35
            for it in items:
                self._rect(s, x + 0.05, y + 0.1, 0.13, 0.13, col)
                self._text(s, x + 0.35, y, 4.9, 0.9,
                           [{"segs": [{"t": it, "sz": 15, "c": self.t.ink,
                                       "font": self.t.fontBody}], "space": 1.18}])
                y += 0.72
        if pageno:
            self._pageno(s, pageno)
        return s

    def closing(self, img, title, lines, footer):
        s = self._slide()
        self._img(s, img, 0, 0, PW, PH)
        self._rect(s, 0, 0, PW, PH, self.t.ink, alpha=52)
        self._text(s, 1.1, 2.1, 11, 1,
                   [{"segs": [{"t": title, "sz": 46, "c": self.t.paper, "b": True,
                               "font": self.t.fontHead}]}])
        self._line(s, 1.13, 3.25, 0.9, self.t.accent, weight=3)
        runs = []
        for ln in lines:
            runs.append({"segs": [{"t": ln, "sz": 18, "c": self.t.paper,
                                   "font": self.t.fontBody}], "space": 1.3, "after": 6})
        self._text(s, 1.1, 3.6, 10.5, 2.5, runs)
        self._text(s, 1.1, PH - 0.8, 11, 0.4,
                   [{"segs": [{"t": footer, "sz": 12, "c": self.t.paper,
                               "font": self.t.fontEN, "spacing": 3}]}])
        return s

    def save(self, path):
        self.prs.save(path)
        return path


THEMES = {
    "spring": Theme(ink="#1d2a24", paper="#f7f4ec", accent="#c0392b", accent2="#2e7d5b",
                    muted="#5a6b62", fontEN="Georgia", fontHead="微软雅黑", fontBody="微软雅黑"),
    "poetry": Theme(ink="#20242e", paper="#f4f1ea", accent="#b8894b", accent2="#3a5a7a",
                    muted="#5c626e", fontEN="Georgia", fontHead="楷体", fontBody="微软雅黑"),
    "seasons": Theme(ink="#12283a", paper="#f5f7f8", accent="#e8843c", accent2="#2f9e9e",
                     muted="#5a6b78", fontEN="Georgia", fontHead="Georgia", fontBody="Segoe UI"),
    "cosmos": Theme(ink="#0f1626", paper="#f2f4fa", accent="#f0a830", accent2="#6c8cff",
                    muted="#5b6478", fontEN="Georgia", fontHead="Georgia", fontBody="Segoe UI"),
}
