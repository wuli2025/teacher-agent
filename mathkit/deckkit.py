# -*- coding: utf-8 -*-
"""数学名师课件共用工具库：统一视觉主题 + PPT 母版 + 配图风格。
所有讲次必须 import 本模块，不得自定义配色/版式，保证 10 讲风格统一。
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------- 视觉主题（深靛蓝学术风 + 琥珀强调） ----------------
INK      = RGBColor(0x11, 0x1B, 0x2E)   # 主文字/深底
NAVY     = RGBColor(0x1B, 0x2A, 0x4A)   # 标题底
SLATE    = RGBColor(0x4A, 0x55, 0x68)   # 次要文字
PAPER    = RGBColor(0xF7, 0xF6, 0xF2)   # 页面底色（暖白）
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0xC8, 0x7E, 0x2A)   # 琥珀，强调
ACCENT2  = RGBColor(0x2E, 0x74, 0x8A)   # 青蓝，第二色
GREEN    = RGBColor(0x3F, 0x7A, 0x5A)
RED      = RGBColor(0xB0, 0x3A, 0x2E)
RULE     = RGBColor(0xD8, 0xD3, 0xC7)   # 分隔线

HEX = lambda c: "#%02X%02X%02X" % (c[0], c[1], c[2])
M_INK, M_ACC, M_ACC2, M_GRN, M_RED = HEX(INK), HEX(ACCENT), HEX(ACCENT2), HEX(GREEN), HEX(RED)
M_PAPER, M_RULE, M_SLATE = HEX(PAPER), HEX(RULE), HEX(SLATE)

CN = "Noto Serif CJK SC"      # 正文中文（衬线，学术感）
CN_SANS = "Noto Sans CJK SC"

rcParams["font.family"] = ["Noto Serif CJK SC", "DejaVu Sans"]
rcParams["axes.unicode_minus"] = False
rcParams["mathtext.fontset"] = "cm"
rcParams["figure.facecolor"] = M_PAPER
rcParams["axes.facecolor"] = M_PAPER
rcParams["savefig.facecolor"] = M_PAPER

W, H = Inches(13.333), Inches(7.5)   # 16:9

# ---------------- 字号下限（硬性） ----------------
# PPT 上任何文字不得小于 20pt。
PPT_MIN_PT = 20
# 图内文字：PNG 贴进 PPT 时通常按 9in 图 → 5.4in 栏缩放（约 0.6×），
# 故图内字号须放大 1/0.6 才能保证贴到幻灯片上 ≥20pt。
# 图内文字：PNG 贴进 PPT 时若被缩小，字也跟着缩小。
# 因此铁律是「不缩小」——贴图宽度必须 >= 画布宽度(英寸)，则图内 20pt = 幻灯片上 20pt。
# picture()/full_picture() 会自动校验，违反即报错，不静默降级。
FIG_MIN_PT = PPT_MIN_PT   # = 20

rcParams["font.size"] = FIG_MIN_PT
rcParams["axes.titlesize"] = FIG_MIN_PT + 4
rcParams["axes.labelsize"] = FIG_MIN_PT
rcParams["xtick.labelsize"] = FIG_MIN_PT
rcParams["ytick.labelsize"] = FIG_MIN_PT
rcParams["legend.fontsize"] = FIG_MIN_PT
rcParams["figure.titlesize"] = FIG_MIN_PT + 6


def _floor(v, lo=FIG_MIN_PT):
    """把字号抬到下限；支持数值与 'small'/'large' 之类的字符串（字符串一律换成下限）。"""
    try:
        return max(float(v), float(lo))
    except (TypeError, ValueError):
        return float(lo)


def _patch_fontsize(cls, names, kw="fontsize"):
    """给 matplotlib 的绘图方法套一层壳，强制 fontsize 不低于 FIG_MIN_PT。
    这样各讲 build 脚本里写死的 fontsize=11/13 也会被自动抬到下限。"""
    for n in names:
        orig = getattr(cls, n, None)
        if orig is None or getattr(orig, "_size_floored", False):
            continue

        def make(f, key):
            def wrapper(*a, **kw_):
                if key in kw_:
                    kw_[key] = _floor(kw_[key])
                elif "size" in kw_:
                    kw_["size"] = _floor(kw_["size"])
                return f(*a, **kw_)
            wrapper._size_floored = True
            return wrapper
        setattr(cls, n, make(orig, kw))


from matplotlib.axes import Axes as _Axes
from matplotlib.figure import Figure as _Figure

_patch_fontsize(_Axes, ["text", "annotate", "set_title", "set_xlabel", "set_ylabel"])
_patch_fontsize(_Figure, ["text", "suptitle"])

_orig_legend = _Axes.legend
if not getattr(_orig_legend, "_size_floored", False):
    def _legend(self, *a, **kw):
        kw["fontsize"] = _floor(kw.get("fontsize", FIG_MIN_PT))
        return _orig_legend(self, *a, **kw)
    _legend._size_floored = True
    _Axes.legend = _legend

_orig_tick = _Axes.tick_params
if not getattr(_orig_tick, "_size_floored", False):
    def _tick_params(self, *a, **kw):
        kw["labelsize"] = _floor(kw.get("labelsize", FIG_MIN_PT))
        return _orig_tick(self, *a, **kw)
    _tick_params._size_floored = True
    _Axes.tick_params = _tick_params

_orig_setxt, _orig_setyt = _Axes.set_xticklabels, _Axes.set_yticklabels
if not getattr(_orig_setxt, "_size_floored", False):
    def _mk_ticklabels(f):
        def w(self, *a, **kw):
            kw["fontsize"] = _floor(kw.get("fontsize", FIG_MIN_PT))
            return f(self, *a, **kw)
        w._size_floored = True
        return w
    _Axes.set_xticklabels = _mk_ticklabels(_orig_setxt)
    _Axes.set_yticklabels = _mk_ticklabels(_orig_setyt)



def new_fig(w=5.4, h=3.3):
    """标准配图画布。w 单位=英寸，且必须 <= 该图在幻灯片上的展示宽度，否则字会被缩小。
    右栏图用默认 (5.4, 3.3)；整幅图用 (9.6, 5.0)。"""
    fig, ax = plt.subplots(figsize=(w, h))
    return fig, ax


def style_axes(ax, xlabel="x", ylabel="y", origin=True, grid=True):
    """统一坐标轴风格：十字坐标轴 + 淡网格。"""
    if grid:
        ax.grid(True, color=M_RULE, lw=0.8, alpha=0.9, zorder=0)
    if origin:
        for s in ("left", "bottom"):
            ax.spines[s].set_position("zero")
            ax.spines[s].set_color(M_INK)
            ax.spines[s].set_linewidth(1.4)
        for s in ("right", "top"):
            ax.spines[s].set_visible(False)
        ax.plot(1, 0, ">k", transform=ax.get_yaxis_transform(), clip_on=False, ms=6)
        ax.plot(0, 1, "^k", transform=ax.get_xaxis_transform(), clip_on=False, ms=6)
    else:
        for s in ("right", "top"):
            ax.spines[s].set_visible(False)
        for s in ("left", "bottom"):
            ax.spines[s].set_color(M_SLATE)
    ax.set_xlabel(xlabel, fontsize=FIG_MIN_PT, color=M_INK, loc="right")
    ax.set_ylabel(ylabel, fontsize=FIG_MIN_PT, color=M_INK, loc="top", rotation=0)
    ax.tick_params(colors=M_SLATE, labelsize=FIG_MIN_PT)


def save_fig(fig, path, transparent=False):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    fig.tight_layout(pad=0.6)
    fig.savefig(path, dpi=200, transparent=transparent)
    plt.close(fig)
    return path


# ---------------- PPT 构件 ----------------
def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, prs, color=PAPER):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def _tb(slide, x, y, w, h, text, size=20, color=INK, bold=False, font=CN,
        align=PP_ALIGN.LEFT, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        run = p.add_run(); run.text = ln
        f = run.font
        f.size = Pt(max(size, PPT_MIN_PT)); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def title_slide(prs, title, subtitle, lecturer="", meta=""):
    s = _blank(prs); _bg(s, prs, INK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), Inches(0.16), Inches(1.9))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); band.shadow.inherit = False
    _tb(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.3), title, 44, WHITE, True)
    _tb(s, Inches(0.95), Inches(3.7), Inches(11.0), Inches(0.9), subtitle, 20, RGBColor(0xC9, 0xCF, 0xDB))
    _tb(s, Inches(0.95), Inches(6.3), Inches(11.0), Inches(0.6),
        "  ·  ".join([x for x in (lecturer, meta) if x]), 14, RGBColor(0x8A, 0x93, 0xA6))
    return s


def section_slide(prs, kicker, title, mins=""):
    """章节过场页（深底），kicker 如 '第 2 幕 · 12 min'。"""
    s = _blank(prs); _bg(s, prs, NAVY)
    _tb(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.5),
        kicker + (f"　|　{mins}" if mins else ""), 16, ACCENT, True, CN_SANS)
    _tb(s, Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.4), title, 40, WHITE, True)
    return s


def content_slide(prs, title, tag=""):
    """内容页骨架：返回 slide，正文区从 y=1.5in 开始。"""
    s = _blank(prs); _bg(s, prs)
    _tb(s, Inches(0.8), Inches(0.45), Inches(10.4), Inches(0.8), title, 30, INK, True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25), Inches(11.7), Emu(12700 * 2))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    if tag:
        tw = max(1.6, 0.34 * len(tag) + 0.5)
        t = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(12.53 - tw), Inches(0.45), Inches(tw), Inches(0.55))
        t.fill.solid(); t.fill.fore_color.rgb = ACCENT
        t.line.fill.background(); t.shadow.inherit = False
        tf = t.text_frame; tf.text = tag
        r = tf.paragraphs[0].runs[0]; r.font.size = Pt(PPT_MIN_PT); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = CN_SANS
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return s


def bullets(slide, items, x=0.85, y=1.6, w=6.2, size=18, color=INK):
    """项目符号列表，items 为 str 或 (文本, 层级) 元组。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        txt, lv = (it, 0) if isinstance(it, str) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.35; p.space_after = Pt(10)
        sz = Pt(max(size if lv == 0 else size - 2, PPT_MIN_PT))
        b = p.add_run(); b.text = ("●  " if lv == 0 else "—  ")
        b.font.size = sz
        b.font.color.rgb = ACCENT if lv == 0 else RULE
        b.font.name = CN_SANS
        r = p.add_run(); r.text = txt
        r.font.size = sz
        r.font.color.rgb = color if lv == 0 else SLATE
        r.font.name = CN
        p.level = lv
    return tb


def callout(slide, text, x=0.85, y=5.6, w=6.2, h=1.2, kind="key"):
    """强调框：kind = key(琥珀) / warn(红) / note(青)。"""
    c = {"key": ACCENT, "warn": RED, "note": ACCENT2}[kind]
    # 20pt 正文一行约需 0.42in；按行数（含折行）撑高，避免文字溢出框外
    n_lines = sum(max(1, -(-len(ln) // int(w * 3.4))) for ln in text.split("\n"))
    h = max(h, 0.40 * n_lines + 0.42)
    if y + h > 7.15:            # 撑高后不得越出下边缘，则整体上移
        y = max(1.45, 7.15 - h)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = c; box.line.width = Pt(1.5); box.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = c
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.15)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run(); r.text = ln
        r.font.size = Pt(PPT_MIN_PT); r.font.color.rgb = INK; r.font.name = CN
    return box


def _assert_not_shrunk(img, w):
    """贴图不得缩小，否则图内字号会跌破 20pt。违反即报错。"""
    from PIL import Image
    with Image.open(img) as im:
        dpi = (im.info.get("dpi") or (200, 200))[0] or 200
        src_in = im.size[0] / float(dpi)
    if w < src_in - 0.02:
        raise ValueError(
            f"{os.path.basename(img)} 画布 {src_in:.2f}in 却只贴 {w:.2f}in，"
            f"缩小 {w/src_in:.0%} → 图内字会小于 20pt。请把 new_fig 宽度改为 <= {w:.2f}")


def picture(slide, img, x=7.3, y=1.55, w=5.4):
    """插图（右栏默认位置）。图片自带浅底，直接贴。"""
    _assert_not_shrunk(img, w)
    return slide.shapes.add_picture(img, Inches(x), Inches(y), width=Inches(w))


def full_picture(slide, img, y=1.5, w=9.6):
    """整幅居中插图。"""
    _assert_not_shrunk(img, w)
    return slide.shapes.add_picture(img, Inches((13.333 - w) / 2), Inches(y), width=Inches(w))


def formula(slide, tex, x=0.85, y=1.7, w=11.6, size=1.0, out=None, color=M_INK):
    """用 matplotlib 渲染 LaTeX 公式为透明 PNG 后贴入（PPT 无原生公式渲染）。
    tex 例：r'$f'(x_0)=\\lim_{\\Delta x\\to 0}\\frac{\\Delta y}{\\Delta x}$'

    关键：用 bbox_inches='tight' 紧致裁剪，PNG 恰好包住整条公式，**绝不裁掉**首尾字符。
    渲染后测量公式真实宽度：
      · 不超过栏宽 w → 按真实宽度贴入，并在 [x, x+w] 内水平居中；
      · 超过栏宽 w   → 整体等比缩到 w（字会略小于设定，但保证完整可读，好过被裁掉）。
    返回贴入的图片 shape，调用方可读 .top/.height 以避免与其他元素穿模。"""
    import PIL.Image
    DPI = 220
    pt = max(30 * size, PPT_MIN_PT + 6)
    fig = plt.figure(figsize=(0.1, 0.1))
    fig.text(0.0, 0.0, tex, ha="left", va="bottom", fontsize=pt, color=color)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    # bbox_inches='tight' 让画布收缩到公式本身，首尾不再被裁
    fig.savefig(out, dpi=DPI, transparent=True, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    with PIL.Image.open(out) as im:
        nat_w = im.size[0] / float(DPI)          # 公式真实宽度（英寸）
        nat_h = im.size[1] / float(DPI)
    place_w = min(w, nat_w)                       # 超宽则缩到栏宽，否则用真实宽度
    place_h = place_w / nat_w * nat_h
    x_off = x + max(0.0, (w - place_w) / 2.0)     # 在栏内水平居中
    return slide.shapes.add_picture(out, Inches(x_off), Inches(y),
                                    width=Inches(place_w), height=Inches(place_h))


def save(prs, path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    try:
        prs.save(path)
    except PermissionError:
        # 原文件被 PowerPoint 占用锁定 → 存成修正版，避免整批中断
        alt = path[:-5] + "-修正版.pptx" if path.endswith(".pptx") else path + "-修正版"
        prs.save(alt)
        print("  (原文件被占用) ->", alt)
        return alt
    return path
