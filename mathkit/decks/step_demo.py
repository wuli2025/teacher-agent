# -*- coding: utf-8 -*-
"""课件B：《高中数学解题步骤演示 · 逐步累积版》

对标参考课件《高中数学解题步骤演示_Manim逐步版》的能力：同一题连续多页、题目区固定、
解题过程逐页累积。三道题分别压中 goal 第五节点名的难点：
  题1 数列   —— 求和号 ∑ 的上下限、多层上下标、错位相减的长等式
  题2 立体几何 —— 向量箭头、坐标、分式套根式、夹角余弦
  题3 线性方程组 —— 增广矩阵的行列与竖线、行变换箭头、负号

产出：out/step-demo/高中数学解题步骤演示.pptx
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.abspath(os.path.join(HERE, "..", "texkit")))

from pptx.util import Inches, Pt                                    # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                              # noqa: E402
from theme import (INK, INK_SOFT, MUTED, CYAN, CORAL, AMBER, BG2, RULE, WHITE,  # noqa: E402
                   FONT, SLIDE_W_IN, MARGIN_L, MARGIN_R, CONTENT_W,
                   BODY_TOP, BODY_BOTTOM, PT_BODY, H_INK, H_CORAL)
from slides import Deck                                             # noqa: E402
from tex import TexPool                                             # noqa: E402
from stepdeck import StepProblem                                    # noqa: E402

OUT = os.path.abspath(os.path.join(HERE, "..", "out", "step-demo"))
TEXDIR = os.path.join(OUT, "tex")
FOOTER = "高中数学 · 解题步骤演示 · 逐步累积版"


def problems(pool):
    p1 = StepProblem(
        pool, no=1, kind="数列 · 错位相减",
        title="错位相减法求 ∑ k·2^(k−1)",
        stem="求和：Sn = 1·2⁰ + 2·2¹ + 3·2² + … + n·2ⁿ⁻¹。",
        stem_note="一个等差数列乘一个等比数列，标准手法就是错位相减。",
        steps=[
            {"head": "写成求和号形式", "tex": [r"S_n=\sum_{k=1}^{n}k\cdot 2^{\,k-1}"],
             "note": "识别结构：等差 × 等比。",
             "speak": "先让学生说出通项 a_k = k·2^(k−1) 是「等差乘等比」，"
                      "这是判定使用错位相减的唯一依据。"},
            {"head": "展开并同乘公比 2", "tex": [
                r"S_n=1+2\cdot 2+3\cdot 2^{2}+\cdots+n\cdot 2^{\,n-1}",
                r"2S_n=1\cdot 2+2\cdot 2^{2}+\cdots+(n-1)2^{\,n-1}+n\cdot 2^{\,n}"],
             "note": "两式的同次项要对齐。",
             "speak": "板书时把两行按 2 的幂次上下对齐，学生才看得出「错位」在哪。"},
            {"head": "作差：错位相减", "tex": [
                r"S_n-2S_n=1+2+2^{2}+\cdots+2^{\,n-1}-n\cdot 2^{\,n}"],
             "note": "中间全部塌缩成等比数列。",
             "speak": "关键一步：作差后中间项系数都变成 1，剩下一个等比数列和一个尾项。"},
            {"head": "用等比求和公式", "tex": [
                r"-S_n=\frac{1-2^{\,n}}{1-2}-n\cdot 2^{\,n}=2^{\,n}-1-n\cdot 2^{\,n}"],
             "note": "注意左边是 −Sn，别丢负号。",
             "speak": "最常见的失分点就在这里丢负号，务必让学生把 −Sn 写清楚。"},
            {"head": "结论", "tex": [r"S_n=(n-1)\,2^{\,n}+1"], "final": True,
             "note": "代 n=1 验算：S₁=1 ✓。",
             "speak": "养成习惯：求和结果一定用 n=1、n=2 回代验算。"},
        ])

    p2 = StepProblem(
        pool, no=2, kind="立体几何 · 二面角",
        title="用空间向量求二面角 B–PC–D 的余弦值",
        stem="四棱锥 P-ABCD 中，底面 ABCD 是边长为 2 的正方形，PA ⊥ 底面，PA = 2。求二面角 B–PC–D 的余弦值。",
        stem_note="向量法三步：建系写坐标 → 求两个半平面的法向量 → 算夹角再定号。",
        steps=[
            {"head": "建系并写出各点坐标", "tex": [
                r"A(0,0,0),\ B(2,0,0),\ C(2,2,0),\ D(0,2,0),\ P(0,0,2)"],
             "note": "以 A 为原点，沿 AB、AD、AP 建三轴。",
             "speak": "PA⊥底面且底面是正方形，天然给出两两垂直的三条棱，直接建右手系。"},
            {"head": "求半平面 PBC 的法向量", "tex": [
                r"\vec{PB}=(2,0,-2),\ \vec{BC}=(0,2,0)\ \Rightarrow\ \vec n_1=(1,0,1)"],
             "note": "法向量与平面内两向量都垂直。",
             "speak": "解 n·PB = 0 与 n·BC = 0，取最简整数解即可，不必单位化。"},
            {"head": "求半平面 PDC 的法向量", "tex": [
                r"\vec{PD}=(0,2,-2),\ \vec{DC}=(2,0,0)\ \Rightarrow\ \vec n_2=(0,1,1)"],
             "note": "同法处理另一个半平面。",
             "speak": "提醒对称性：把 x、y 互换即可得到 n₂，可作为检验。"},
            {"head": "算两法向量夹角的余弦", "tex": [
                r"\cos\langle \vec n_1,\vec n_2\rangle=\frac{\vec n_1\cdot\vec n_2}"
                r"{|\vec n_1|\,|\vec n_2|}=\frac{1}{\sqrt2\cdot\sqrt2}=\frac{1}{2}"],
             "note": "这只是法向量夹角，还没定号。",
             "speak": "务必强调：法向量夹角 ≠ 二面角，两者要么相等要么互补。"},
            {"head": "结论：判断钝角并定号", "tex": [
                r"\angle(B\text{-}PC\text{-}D)\ \text{为钝角}\ \Rightarrow\ "
                r"\cos\angle(B\text{-}PC\text{-}D)=-\frac{1}{2}"],
             "final": True, "note": "观察图形定号，别只算不判。",
             "speak": "由图可见二面角张开超过 90°，故取补角，余弦为负，二面角为 120°。"},
        ])

    p3 = StepProblem(
        pool, no=3, kind="线性方程组 · 高斯消元",
        title="用增广矩阵的行变换解三元一次方程组",
        stem="解方程组：x + 2y − z = 2，2x − y + 3z = 9，−x + y + 2z = 3。",
        stem_note="把方程组写成增广矩阵，用行变换化成阶梯形再回代，步骤最不易错。",
        steps=[
            {"head": "写出增广矩阵", "tex": [
                r"\left[\begin{array}{ccc|c}1&2&-1&2\\ 2&-1&3&9\\ -1&1&2&3\end{array}\right]"],
             "note": "竖线左边是系数，右边是常数。",
             "speak": "强调竖线的意义：它把系数矩阵和常数列分开，行变换要整行一起做。"},
            {"head": "第一列消元", "tex": [
                r"\xrightarrow{\;r_2-2r_1,\ r_3+r_1\;}"
                r"\left[\begin{array}{ccc|c}1&2&-1&2\\ 0&-5&5&5\\ 0&3&1&5\end{array}\right]"],
             "note": "用第一行把第一列其余项打成 0。",
             "speak": "行变换写在箭头上方，批改时一眼能看出每一步做了什么。"},
            {"head": "第二列消元，化成阶梯形", "tex": [
                r"\xrightarrow{\;r_2\div(-5),\ r_3-3r_2\;}"
                r"\left[\begin{array}{ccc|c}1&2&-1&2\\ 0&1&-1&-1\\ 0&0&4&8\end{array}\right]"],
             "note": "先把主元化成 1，后面更好算。",
             "speak": "把 r₂ 除以 −5 得到主元 1，是减少分数运算的关键一招。"},
            {"head": "回代求解", "tex": [
                r"4z=8\ \Rightarrow\ z=2,\qquad y-z=-1\ \Rightarrow\ y=1",
                r"x+2y-z=2\ \Rightarrow\ x=2-2\cdot 1+2=2"],
             "note": "从最后一行往上逐个代回。",
             "speak": "回代顺序固定：先 z，再 y，最后 x，一步一个方程。"},
            {"head": "结论", "tex": [r"(x,\,y,\,z)=(2,\,1,\,2)"], "final": True,
             "note": "把解代回三个原方程逐一验证。",
             "speak": "线性方程组必须验算：代回三式全部成立才算做完。"},
        ])
    return [p1, p2, p3]


def build(deck, pool, probs):
    # ── 封面 ──
    s = deck.blank("封面", bg=INK, footer=False)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(2.24),
                              Inches(0.14), Inches(1.64))
    band.fill.solid(); band.fill.fore_color.rgb = AMBER
    band.line.fill.background(); band.shadow.inherit = False
    deck.label(s, MARGIN_L + 0.34, 1.66, 10.6, "高中数学 · 解题步骤演示", pt=18, color=CYAN)
    deck.para(s, MARGIN_L + 0.34, 2.14, 11.2, "三道综合题的逐步拆解", pt=46, color=WHITE,
              bold=True, spacing=1.05, role="h1")
    deck.para(s, MARGIN_L + 0.34, 3.28, 11.2,
              "同一题连续多页：题目区固定不动，解题过程逐页累积一行。",
              pt=22, color=RULE, spacing=1.0)
    deck.rule(s, MARGIN_L + 0.34, 4.10, 3.2, color=AMBER, thick=0.03)
    deck.para(s, MARGIN_L + 0.34, 4.36, 11.2,
              "一句核心理解：讲题的价值不在答案，而在每一步「为什么这样写」。",
              pt=24, color=WHITE, spacing=1.25)
    deck.label(s, MARGIN_L + 0.34, 6.52, 11.2, FOOTER, pt=15, color=MUTED, bold=False)
    deck.notes(s, "开场说明用法：本课件不依赖动画，翻页即推进一步；"
                  "每页左栏是当前步，右栏是已经走过的全部步骤。")

    # ── 目录 ──
    s = deck.blank("目录")
    deck.header(s, "三道题，三种最常考的结构", kicker="课件结构", tag="导航")
    items = [
        ("01", "数列", "错位相减求 ∑ k·2^(k−1)", "求和号上下限、多层上下标、长等式对齐", CYAN),
        ("02", "立体几何", "空间向量求二面角余弦", "向量箭头、坐标、分式套根式、定号", AMBER),
        ("03", "线性方程组", "增广矩阵行变换求解", "矩阵竖线、行变换箭头、负号、回代", CORAL),
    ]
    ch = (BODY_BOTTOM - BODY_TOP - 0.44) / 3
    yy = BODY_TOP
    for no, kind, title, points, col in items:
        deck.card(s, MARGIN_L, yy, CONTENT_W, ch, fill=BG2, line=RULE, bar=col)
        deck.text(s, MARGIN_L + 0.24, yy, 1.3, ch, no, pt=38, color=col, bold=True,
                  spacing=1.0, space_after=0, anchor="ctr", role="h1")
        deck.label(s, MARGIN_L + 1.70, yy + 0.26, 2.6, kind, pt=16, color=col)
        deck.para(s, MARGIN_L + 1.66, yy + 0.62, 4.6, title, pt=22, color=INK, bold=True,
                  spacing=1.05)
        ph = deck.measure_text_h(points, 20, CONTENT_W - 6.9, spacing=1.28, space_after=0)
        deck.para(s, MARGIN_L + 6.6, yy + (ch - ph) / 2, CONTENT_W - 6.85, points,
                  pt=20, color=INK_SOFT, spacing=1.28)
        yy += ch + 0.22
    deck.notes(s, "先给学生看结构：三道题分别练三种「一看就知道该用哪招」的题型。")

    # ── 三道题的逐步页 ──
    for p in probs:
        p.emit(deck)

    # ── 方法总纲 ──
    s = deck.blank("方法总纲")
    deck.header(s, "三题共用的一套做题动作", kicker="方法提炼", tag="总纲", tag_color=CORAL)
    cols = [
        ("识结构", "看到什么就用什么",
         "等差×等比用错位相减；有垂直棱就建系用向量；三元一次上矩阵。"),
        ("列步骤", "先写框架再算数",
         "先按行写下这一题的固定动作，再逐行填数，别边想边算跳步。"),
        ("守细节", "负号、定号、竖线",
         "错位相减守负号；二面角守定号；行变换整行同步，别漏常数列。"),
        ("必验算", "回代是最后一步",
         "求和用 n=1 回代；法向量点乘验零；方程组把解代回三式。"),
    ]
    cw = (CONTENT_W - 0.75) / 4
    ytop = BODY_TOP + 0.16
    bh = max(deck.measure_text_h(b, 20, cw - 0.32, spacing=1.30, space_after=0)
             for _h, _s, b in cols)
    cardh = 1.78 + bh + 0.18
    for i, (h1, sub, body) in enumerate(cols):
        cx = MARGIN_L + i * (cw + 0.25)
        col = [CYAN, AMBER, CORAL, INK][i]
        deck.card(s, cx, ytop, cw, cardh, fill=BG2, line=RULE, bar=col)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.18), Inches(ytop + 0.22),
                                 Inches(0.52), Inches(0.52))
        num.fill.solid(); num.fill.fore_color.rgb = col
        num.line.fill.background(); num.shadow.inherit = False
        num.name = "tag_num"
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf = num.text_frame; tf.word_wrap = False
        p_ = tf.paragraphs[0]; p_.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p_.add_run(); r.text = str(i + 1)
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        deck.para(s, cx + 0.16, ytop + 0.92, cw - 0.32, h1, pt=22, color=INK, bold=True,
                  spacing=1.0)
        deck.label(s, cx + 0.18, ytop + 1.38, cw - 0.36, sub, pt=16, color=col)
        deck.para(s, cx + 0.16, ytop + 1.78, cw - 0.32, body, pt=20, color=INK_SOFT,
                  spacing=1.30)
    cy = ytop + cardh + 0.22
    deck.card(s, MARGIN_L, cy, CONTENT_W, BODY_BOTTOM - cy, fill=WHITE, line=CORAL, bar=CORAL)
    deck.para(s, MARGIN_L + 0.22, cy + 0.20, CONTENT_W - 0.44,
              "离堂检验：任选一题，只写「四个动作」的框架，不算数。",
              pt=22, color=INK, bold=True, spacing=1.0)
    deck.notes(s, "收束：三道题看似无关，动作是同一套。让学生当场默写框架，"
                  "下节课用它讲评作业。")


def main():
    os.makedirs(OUT, exist_ok=True)
    pool = TexPool(TEXDIR, scale=4)
    probs = problems(pool)
    pool.render()
    deck = Deck(footer_left=FOOTER, pool=pool)
    build(deck, pool, probs)
    path = deck.save(os.path.join(OUT, "高中数学解题步骤演示.pptx"))
    print(f"课件B: {len(deck.pages)} 页 -> {path}")
    return path


if __name__ == "__main__":
    main()
