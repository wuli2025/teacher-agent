# -*- coding: utf-8 -*-
"""数学示意图 —— 用 PPT 原生矢量图形按坐标精确绘制。

为什么不交给文生图：实测 MiniMax image-01（开/关 prompt_optimizer 都试过）画不对数学关系——
要求「直线在中点断开、留一个空心点」，出来的是装饰性波浪块面、实心点、甚至冒出类字母笔画，
违反 goal「数学关系正确」与「不得含字母」两条硬要求。
所以本工具箱的分工是：
  · 承载数学关系的示意图 → 本模块矢量绘制（确定性、可编辑、可无损缩放）
  · 不承载数学关系的情境插图 → 文生图（genimg.py）
"""
import os
import sys

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from theme import INK, INK_SOFT, MUTED, CYAN, CORAL, AMBER, RULE, WHITE, BG2


def _freeform(slide, pts, color, width_pt=2.5, dash=False):
    """按点序画折线（点足够密即为光滑曲线）。pts 为英寸坐标。"""
    b = slide.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]))
    b.add_line_segments([(Inches(px), Inches(py)) for px, py in pts[1:]], close=False)
    sh = b.convert_to_shape()
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(width_pt)
    sh.shadow.inherit = False
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return sh


def _arrow(slide, x1, y1, x2, y2, color=INK_SOFT, width_pt=1.5):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    ln = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return c


def _dash(slide, x1, y1, x2, y2, color=MUTED):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(1.1)
    c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def _open_dot(slide, cx, cy, r=0.085, color=CORAL):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                               Inches(2 * r), Inches(2 * r))
    d.fill.solid(); d.fill.fore_color.rgb = WHITE
    d.line.color.rgb = color; d.line.width = Pt(2.2)
    d.shadow.inherit = False
    d.name = "dia_opendot"
    return d


def _solid_dot(slide, cx, cy, r=0.06, color=CORAL):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                               Inches(2 * r), Inches(2 * r))
    d.fill.solid(); d.fill.fore_color.rgb = color
    d.line.fill.background(); d.shadow.inherit = False
    d.name = "dia_dot"
    return d


def _frame(slide, x, y, w, h):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.color.rgb = RULE; r.line.width = Pt(1.0)
    r.shadow.inherit = False
    r.name = "dia_frame"
    try:
        r.adjustments[0] = 0.05
    except Exception:
        pass
    return r


# ────────────────────────── 图 1：可去间断点（洞）──────────────────────────
def hole_diagram(slide, x, y, w, h, hole_t=0.5):
    """y = x + 1 型直线，在 x = 1 处挖掉一个点：一条直线 + 一个空心点 + 两条虚线引导。

    (x, y, w, h) 为绘图区外框（英寸）。图形自适应该区域，等比无关（纯矢量）。
    """
    _frame(slide, x, y, w, h)
    pad = 0.30
    ox, oy = x + pad + 0.18, y + h - pad - 0.10          # 原点
    ax_w, ax_h = w - 2 * pad - 0.28, h - 2 * pad - 0.05  # 轴长
    _arrow(slide, ox - 0.12, oy, ox + ax_w, oy)          # x 轴
    _arrow(slide, ox, oy + 0.12, ox, oy - ax_h)          # y 轴

    # 直线：从左下到右上，占轴内 10%~92%
    x0, x1 = ox + 0.10 * ax_w, ox + 0.92 * ax_w
    y0, y1 = oy - 0.16 * ax_h, oy - 0.88 * ax_h
    hx = x0 + hole_t * (x1 - x0)
    hy = y0 + hole_t * (y1 - y0)
    gap = 0.075                                          # 断口半长（沿线方向）
    dx, dy = (x1 - x0), (y1 - y0)
    L = (dx * dx + dy * dy) ** 0.5
    ux, uy = dx / L, dy / L
    _freeform(slide, [(x0, y0), (hx - gap * ux, hy - gap * uy)], INK, 3.0)
    _freeform(slide, [(hx + gap * ux, hy + gap * uy), (x1, y1)], INK, 3.0)

    _dash(slide, ox, hy, hx, hy)
    _dash(slide, hx, oy, hx, hy)
    _open_dot(slide, hx, hy)
    return (hx, hy)


# ────────────────────────── 图 2：夹逼定理 ──────────────────────────
def squeeze_diagram(slide, x, y, w, h):
    """上下两条曲线向右收拢，把中间的曲线夹向同一个极限值。"""
    _frame(slide, x, y, w, h)
    pad = 0.30
    ox, oy = x + pad + 0.18, y + h - pad - 0.10
    ax_w, ax_h = w - 2 * pad - 0.28, h - 2 * pad - 0.05
    _arrow(slide, ox - 0.12, oy, ox + ax_w, oy)
    _arrow(slide, ox, oy + 0.12, ox, oy - ax_h)

    xa, xb = ox + 0.08 * ax_w, ox + 0.90 * ax_w
    Lv = oy - 0.55 * ax_h                                # 共同极限所在高度
    amp = 0.30 * ax_h                                    # 初始张口半宽
    N = 40
    up, lo, mid = [], [], []
    for i in range(N + 1):
        t = i / N
        px = xa + t * (xb - xa)
        k = (1 - t) ** 1.6                                # 张口随 x 收拢
        up.append((px, Lv - amp * k))
        lo.append((px, Lv + amp * k))
        # 中间曲线在上下之间来回摆动，被越夹越紧
        import math
        mid.append((px, Lv - amp * k * 0.72 * math.cos(6.0 * t)))
    _freeform(slide, up, INK, 2.6)
    _freeform(slide, lo, CYAN, 2.6)
    _freeform(slide, mid, CORAL, 3.0)
    _dash(slide, ox, Lv, xb + 0.06, Lv)
    _solid_dot(slide, xb, Lv)
    return (xb, Lv)
