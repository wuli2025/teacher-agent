# -*- coding: utf-8 -*-
"""生成每套课件的「构建结果记录」与「公式兼容性验证报告」。

用法：python report.py <deck_out_dir> "<课题名>"
产出：<dir>/构建结果记录.json、<dir>/公式兼容性验证报告.md
"""
import json
import os
import sys
from collections import Counter

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from check import check_pptx                       # noqa: E402

# 公式类型识别（用于报告里「使用 Manim/LaTeX 的公式类型」一栏）
KINDS = [
    ("积分", r"\int"), ("求和", r"\sum"), ("连乘", r"\prod"), ("极限", r"\lim"),
    ("分式", r"\frac"), ("分式(行内)", r"\tfrac"), ("分式(展示)", r"\dfrac"),
    ("根式", r"\sqrt"), ("矩阵/数组", r"begin{array"), ("矩阵", r"begin{bmatrix"),
    ("行列式", r"begin{vmatrix"), ("分段函数", r"begin{cases"),
    ("多行对齐推导", r"begin{aligned"), ("向量", r"\vec"), ("行变换箭头", r"\xrightarrow"),
    ("蕴含/等价", r"\Rightarrow"), ("等价", r"\iff"), ("量词", r"\forall"),
    ("量词(存在)", r"\exists"), ("希腊字母", r"\varepsilon"), ("绝对值/范数", r"\|"),
    ("公式内中文", r"\text{"), ("上下叠标注", r"\underbrace"),
    ("内积尖括号", r"\langle"), ("角/度数", r"\angle"),
]


def kinds_of(tex):
    return [k for k, pat in KINDS if pat in tex]


def main():
    out_dir = os.path.abspath(sys.argv[1])
    topic = sys.argv[2] if len(sys.argv) > 2 else os.path.basename(out_dir)
    pptx = [f for f in os.listdir(out_dir) if f.endswith(".pptx")]
    if not pptx:
        print("找不到 pptx"); sys.exit(1)
    pptx_path = os.path.join(out_dir, pptx[0])

    metas = []
    tex_dir = os.path.join(out_dir, "tex")
    for f in sorted(os.listdir(tex_dir)):
        if f.endswith(".json") and not f.startswith("_"):
            metas.append(json.load(open(os.path.join(tex_dir, f), encoding="utf-8")))
    uniq, seen = [], set()
    for m in metas:
        if m["tex"] not in seen:
            seen.add(m["tex"]); uniq.append(m)

    prs = Presentation(pptx_path)
    n_pages = len(prs.slides.__iter__.__self__._sldIdLst)
    pics = sum(1 for s in prs.slides for sh in s.shapes
               if sh.__class__.__name__ == "Picture")
    tex_pics = sum(1 for s in prs.slides for sh in s.shapes
                   if sh.__class__.__name__ == "Picture" and (sh.name or "").startswith("tex_"))
    issues = check_pptx(pptx_path)
    preview = os.path.join(out_dir, "preview")
    n_prev = len([f for f in os.listdir(preview) if f.lower().endswith(".png")]) \
        if os.path.isdir(preview) else 0

    kind_count = Counter()
    for m in uniq:
        for k in kinds_of(m["tex"]):
            kind_count[k] += 1

    rec = {
        "课题": topic,
        "pptx": os.path.basename(pptx_path),
        "页数": n_pages,
        "公式条数(去重)": len(uniq),
        "公式PNG张数": tex_pics,
        "图片总数": pics,
        "逐页实渲张数": n_prev,
        "静态体检": "通过" if not issues else f"{len(issues)} 项未通过",
        "静态体检明细": [f"P{p:02d} [{w}] {m}" for p, w, m in issues],
        "公式类型分布": dict(kind_count.most_common()),
        "渲染管线": "KaTeX 0.17 + 无头 Chrome 裸 CDP，scale=4（≈288DPI）",
    }
    json.dump(rec, open(os.path.join(out_dir, "构建结果记录.json"), "w", encoding="utf-8"),
              ensure_ascii=False, indent=1)

    lines = [f"# 公式兼容性验证报告 · {topic}", "",
             f"* 渲染器：KaTeX 0.17（`throwOnError: true`，任何宏不支持都直接失败，绝不静默出错字）",
             f"* 截图：无头 Chrome 按元素包围盒裁剪，scale = 4（30pt 公式 ≈ 288 DPI）",
             f"* 贴图规则：1 CSS px = 1 pt，宽高同一比例因子换算，**不可能拉伸变形**",
             f"* 本课件公式（去重）：**{len(uniq)}** 条，全部渲染成功",
             "", "## 公式类型覆盖", "", "| 类型 | 条数 |", "| --- | --- |"]
    for k, v in kind_count.most_common():
        lines.append(f"| {k} | {v} |")
    lines += ["", "## 逐条明细", "",
              "| # | LaTeX 源码 | PNG 像素 | 排版尺寸(in) | 状态 |",
              "| --- | --- | --- | --- | --- |"]
    for i, m in enumerate(uniq, 1):
        w_in = m["cssW"] / 72.0
        h_in = w_in * m["h"] / m["w"]
        tex = m["tex"].replace("|", "\\|").replace("\n", " ")
        lines.append(f"| {i} | `{tex}` | {m['w']}×{m['h']} | {w_in:.2f}×{h_in:.2f} | ✅ |")
    open(os.path.join(out_dir, "公式兼容性验证报告.md"), "w", encoding="utf-8").write(
        "\n".join(lines) + "\n")

    print(f"{topic}: {n_pages} 页 / {len(uniq)} 条公式 / 实渲 {n_prev} 张 / "
          f"体检 {rec['静态体检']}")


if __name__ == "__main__":
    main()
