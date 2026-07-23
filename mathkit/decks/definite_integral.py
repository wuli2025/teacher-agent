# -*- coding: utf-8 -*-
"""课件C：《定积分与微积分基本定理 —— 从和式极限到牛顿-莱布尼茨公式》

内容来源：mathkit/build_02.py（第02讲）。原课件只作知识来源，不被修改。
选题理由：对十讲做公式复杂度审查，第02讲仅次于第10讲（积分 13、求和 4、多层上下标 25），
且难点集中在「积分号与上下限」「求和号与上下限」「长分式」—— 正好是课件A/B 未覆盖的一组。

产出：out/definite-integral/ 下的 pptx + 公式PNG + 情境图 + Manim 源码 + 逐页实渲 + 报告
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.abspath(os.path.join(HERE, "..", "texkit")))

from pptx.util import Inches, Pt                                    # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                              # noqa: E402
from theme import (INK, INK_SOFT, MUTED, CYAN, CORAL, AMBER, BG2, RULE, WHITE,  # noqa: E402
                   FONT, SLIDE_W_IN, MARGIN_L, MARGIN_R, CONTENT_W,
                   BODY_TOP, BODY_BOTTOM, PT_BODY, PT_TEX_DISPLAY,
                   H_INK, H_CYAN, H_CORAL)
from slides import Deck                                             # noqa: E402
from tex import TexPool                                             # noqa: E402
from stepdeck import StepProblem                                    # noqa: E402
from diagram import riemann_diagram                                 # noqa: E402

OUT = os.path.abspath(os.path.join(HERE, "..", "out", "definite-integral"))
IMG = os.path.join(OUT, "img")
TEXDIR = os.path.join(OUT, "tex")
FOOTER = "高中数学 · 定积分与微积分基本定理 · 和式极限与 N-L 公式"

# 承载数学关系的图走 diagram.py 矢量绘制；文生图只做封面的情境意象。
IMAGES = {
    "slices": ("deep indigo and teal abstract illustration, a solid curved shape built up "
               "from a stack of many extremely thin horizontal slices, the slices getting "
               "thinner and denser toward the top until they merge into a smooth surface, "
               "dark background, soft rim light, cinematic minimal, "
               "no people, no text, no letters, no numbers"),
}


def register(pool):
    T = {}
    A = lambda k, tex, pt=PT_TEX_DISPLAY, c=H_INK: T.__setitem__(k, pool.add(tex, pt=pt, color=c))

    A("q_area", r"S=\int_{0}^{1}x^{2}\,\mathrm{d}x=\ ?", 32, H_CORAL)
    A("q_rect", r"S_{\text{矩形}}=\text{底}\times\text{高}", 26)
    # 核心公式：定积分定义
    A("define",
      r"\int_{a}^{b} f(x)\,\mathrm{d}x=\lim_{n\to\infty}\sum_{i=1}^{n}"
      r"f(\xi_i)\cdot\frac{b-a}{n}", 30)
    A("k1", r"\Delta x=\frac{b-a}{n}", 28, H_CYAN)
    A("k2", r"f(\xi_i)\,\Delta x", 28, H_CYAN)
    # \textstyle 让 lim/∑ 的上下限改排到右侧：展示式样会叠成三层，小卡片里放不下
    A("k3", r"\textstyle\lim_{n\to\infty}\sum_{i=1}^{n}", 26, H_CYAN)
    # 牛顿-莱布尼茨
    A("nl", r"\int_{a}^{b} f(x)\,\mathrm{d}x=F(b)-F(a)=\bigl[F(x)\bigr]_{a}^{b},"
            r"\quad F'(x)=f(x)", 27)
    A("area_fn", r"A(x)=\int_{a}^{x} f(t)\,\mathrm{d}t\ \Longrightarrow\ A'(x)=f(x)", 24)
    # 迁移
    A("move", r"s=\int_{t_1}^{t_2} v(t)\,\mathrm{d}t,\qquad "
              r"\ell=\int_{t_1}^{t_2}\bigl|v(t)\bigr|\,\mathrm{d}t", 24)
    A("work", r"W=\int_{a}^{b} F(x)\,\mathrm{d}x", 26, H_CORAL)
    # 易错
    A("err1", r"\int_{b}^{a} f=-\int_{a}^{b} f", 22, H_CORAL)
    A("err2", r"\int_{0}^{2}\bigl|x^{2}-1\bigr|\mathrm{d}x\neq"
              r"\left[\tfrac{x^{3}}{3}-x\right]_{0}^{2}", 21, H_CORAL)
    A("err3", r"S=\int_{a}^{b}\bigl|f(x)\bigr|\,\mathrm{d}x", 22, H_CORAL)
    # 练习
    A("ex1", r"\int_{0}^{1}\left(3x^{2}+2x\right)\mathrm{d}x", 26)
    A("ex2", r"v(t)=t^{2}-4t+3\ (0\le t\le 3)", 24)
    # 离堂
    A("out1", r"\int_{0}^{2}(x-1)\,\mathrm{d}x", 28)
    return T


def problems(pool):
    p1 = StepProblem(
        pool, no=1, kind="用定义计算定积分",
        title="不用任何公式，硬算 ∫₀¹ x² dx",
        stem="用定积分的定义（分割 → 近似代替 → 求和 → 取极限）计算 ∫₀¹ x² dx。",
        stem_note="定义题的价值不在结果，而在于看清「和式的极限」这四个字。",
        steps=[
            {"head": "分割：等分区间", "tex": [
                r"\Delta x=\frac{1}{n},\qquad \xi_i=\frac{i}{n}\quad(i=1,2,\dots,n)"],
             "note": "把 [0,1] 等分成 n 份",
             "speak": "取右端点作 ξᵢ 只是为了算式好看；取左端点、取中点，极限都一样，"
                      "这正是定积分「与取法无关」的深意。"},
            {"head": "近似代替并求和", "tex": [
                r"S_n=\sum_{i=1}^{n}\left(\frac{i}{n}\right)^{2}\cdot\frac{1}{n}"
                r"=\frac{1}{n^{3}}\sum_{i=1}^{n} i^{2}"],
             "note": "每条矩形：高 × 宽",
             "speak": "先把 1/n³ 提到求和号外面——这一步是能不能算下去的关键。"},
            {"head": "代平方和公式", "tex": [
                r"=\frac{1}{n^{3}}\cdot\frac{n(n+1)(2n+1)}{6}"
                r"=\frac{1}{6}\left(1+\frac{1}{n}\right)\left(2+\frac{1}{n}\right)"],
             "note": "1²+2²+…+n² 的封闭式",
             "speak": "平方和公式要背熟；化成含 1/n 的形式，才好一眼看出极限。"},
            {"head": "取极限", "tex": [
                r"\lim_{n\to\infty}S_n=\frac{1}{6}\cdot 1\cdot 2=\frac{1}{3}"],
             "note": "n→∞ 时 1/n→0",
             "speak": "到这一步才第一次用到「极限」；前三步都只是代数运算。"},
            {"head": "结论", "tex": [r"\int_{0}^{1}x^{2}\,\mathrm{d}x=\frac{1}{3}"],
             "final": True, "note": "和式的极限就是定积分",
             "speak": "对照后面的 N-L 公式：[x³/3]₀¹ = 1/3，一行就出来了——"
                      "这就是基本定理的威力，但今天先让学生体会它有多难得。"},
        ])

    p2 = StepProblem(
        pool, no=2, kind="N-L 公式 · 含绝对值",
        title="用牛顿-莱布尼茨公式算 ∫₀² |x²−1| dx",
        stem="计算 ∫₀² |x²−1| dx。",
        stem_note="带绝对值的定积分，必须先按被积函数的正负拆区间，再分段用 N-L 公式。",
        steps=[
            {"head": "先去绝对值：判正负", "tex": [
                r"\bigl|x^{2}-1\bigr|=\begin{cases}1-x^{2}, & 0\le x\le 1\\[2pt]"
                r"x^{2}-1, & 1<x\le 2\end{cases}"],
             "note": "分界点是 x = 1",
             "speak": "让学生先解 x²−1=0 找到分界点，再判断每段的正负号，不要凭感觉。"},
            {"head": "按分界点拆区间", "tex": [
                r"\int_{0}^{2}\bigl|x^{2}-1\bigr|\mathrm{d}x"
                r"=\int_{0}^{1}\left(1-x^{2}\right)\mathrm{d}x"
                r"+\int_{1}^{2}\left(x^{2}-1\right)\mathrm{d}x"],
             "note": "区间可加性",
             "speak": "定积分对区间可加，这是拆的依据；拆完每段的被积函数都不再带绝对值。"},
            {"head": "各段求原函数", "tex": [
                r"F_1(x)=x-\frac{x^{3}}{3},\qquad F_2(x)=\frac{x^{3}}{3}-x"],
             "note": "求导验回去就对",
             "speak": "原函数不唯一（差常数），但代入上下限时常数会抵消，取最简即可。"},
            {"head": "代入上下限", "tex": [
                r"=\left[x-\frac{x^{3}}{3}\right]_{0}^{1}"
                r"+\left[\frac{x^{3}}{3}-x\right]_{1}^{2}"
                r"=\frac{2}{3}+\left(\frac{2}{3}+\frac{2}{3}\right)"],
             "note": "上限减下限，别写反",
             "speak": "方括号加上下标是规范写法，批卷时这一步不写会扣过程分。"},
            {"head": "结论", "tex": [r"\int_{0}^{2}\bigl|x^{2}-1\bigr|\mathrm{d}x=2"],
             "final": True, "note": "它同时也是面积",
             "speak": "因为被积函数取了绝对值，这个值就等于曲线与 x 轴围成的面积。"},
        ])
    return [p1, p2]


def build(deck, pool, T, probs):
    # ── P1 封面 ──
    s = deck.blank("封面", bg=INK, footer=False)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(2.30),
                              Inches(0.14), Inches(1.62))
    band.fill.solid(); band.fill.fore_color.rgb = AMBER
    band.line.fill.background(); band.shadow.inherit = False
    tw = 7.05
    cover = os.path.join(IMG, "slices.png")
    if os.path.exists(cover):
        deck.picture(s, cover, MARGIN_L + 0.34 + tw + 0.40, 2.16,
                     SLIDE_W_IN - MARGIN_R - (MARGIN_L + 0.34 + tw + 0.40), 2.66, mode="cover")
    deck.label(s, MARGIN_L + 0.34, 1.72, tw, "高中数学 · 选择性必修", pt=18, color=CYAN)
    deck.para(s, MARGIN_L + 0.34, 2.20, tw, "定积分与微积分基本定理", pt=42, color=WHITE,
              bold=True, spacing=1.05, role="h1")
    deck.para(s, MARGIN_L + 0.34, 3.34, tw, "教学主线：直观理解 → 公式表达 → 迁移应用",
              pt=22, color=RULE, spacing=1.0)
    deck.rule(s, MARGIN_L + 0.34, 4.16, 3.2, color=AMBER, thick=0.03)
    deck.para(s, MARGIN_L + 0.34, 4.42, tw,
              "一句核心理解：定积分是「和式的极限」，而 N-L 公式让这个极限变成一次减法。",
              pt=23, color=WHITE, spacing=1.25)
    deck.label(s, MARGIN_L + 0.34, 6.52, 11.2, FOOTER, pt=15, color=MUTED, bold=False)
    deck.notes(s, "开场先问：矩形面积会算，曲边梯形呢？把「化曲为直」这四个字写在黑板角落，"
                  "整节课都在兑现它。")

    # ── P2 问题驱动 ──
    s = deck.blank("问题驱动")
    deck.header(s, "直边会算，曲边呢？", kicker="问题驱动", tag="导入")
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "小学就会算矩形面积；可当上边界换成一条曲线，公式立刻失效：", pt=21)
    y += 0.10
    half = (CONTENT_W - 0.40) / 2
    deck.card(s, MARGIN_L, y, half, 1.30, fill=BG2, line=RULE, bar=CYAN)
    pool.place(s, T["q_rect"], MARGIN_L + 0.14, y, max_w=half - 0.28, max_h=1.30,
               align="ctr", valign="ctr", name="tex_q_rect")
    rx = MARGIN_L + half + 0.40
    deck.card(s, rx, y, half, 1.30, fill=BG2, line=CORAL, bar=CORAL)
    pool.place(s, T["q_area"], rx + 0.14, y, max_w=half - 0.28, max_h=1.30,
               align="ctr", valign="ctr", name="tex_q_area")
    y += 1.30 + 0.24
    cw0 = (CONTENT_W - 0.5) / 3
    cards = [
        ("直边可拆", "矩形、三角形、梯形都能用现成公式拼出来。", CYAN),
        ("曲边卡住", "上边界每一点高度都不同，没有一个「高」可用。", AMBER),
        ("于是化曲为直", "先用无数细矩形近似，再让它们无限变薄。", CORAL),
    ]
    bh = max(deck.measure_text_h(b, 20, cw0 - 0.40, spacing=1.26, space_after=0)
             for _h, b, _c in cards)
    ch = 0.62 + bh + 0.20
    for i, (h1, b1, col) in enumerate(cards):
        cx = MARGIN_L + i * (cw0 + 0.25)
        deck.card(s, cx, y, cw0, ch, fill=WHITE, line=RULE, bar=col)
        deck.para(s, cx + 0.20, y + 0.16, cw0 - 0.40, h1, pt=21, color=col, bold=True, spacing=1.0)
        deck.para(s, cx + 0.20, y + 0.62, cw0 - 0.40, b1, pt=20, color=INK_SOFT, spacing=1.26)
    deck.notes(s, "让学生自己说出「曲边没法直接套公式」，再引出「用直的去逼近曲的」。约 4 分钟。")

    # ── P3 学习路线 ──
    s = deck.blank("学习路线")
    deck.header(s, "本节课的三段路线", kicker="学习路线", tag="路线")
    routes = [
        ("① 直观理解", "分割求和", "把曲边梯形切成细矩形，算出一个和式，再让它无限变细。"),
        ("② 公式表达", "和式的极限", "把这个过程写成定积分的定义，再由基本定理换成 N-L 公式。"),
        ("③ 迁移应用", "面积·位移·功", "同一个「累积」结构，可以算面积、算位移、算变力做功。"),
    ]
    cw = (CONTENT_W - 0.7) / 3
    for i, (tag, head, body) in enumerate(routes):
        cx = MARGIN_L + i * (cw + 0.35)
        col = [CYAN, AMBER, CORAL][i]
        deck.card(s, cx, BODY_TOP + 0.24, cw, 3.36, fill=BG2 if i % 2 == 0 else WHITE,
                  line=RULE, bar=col)
        deck.label(s, cx + 0.22, BODY_TOP + 0.44, cw - 0.44, tag, pt=17, color=col)
        deck.para(s, cx + 0.22, BODY_TOP + 0.88, cw - 0.44, head, pt=28, color=INK,
                  bold=True, spacing=1.05)
        deck.para(s, cx + 0.22, BODY_TOP + 1.62, cw - 0.44, body, pt=20, color=INK_SOFT,
                  spacing=1.30)
        if i < 2:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx + cw + 0.06),
                                    Inches(BODY_TOP + 1.70), Inches(0.23), Inches(0.30))
            ar.fill.solid(); ar.fill.fore_color.rgb = RULE
            ar.line.fill.background(); ar.shadow.inherit = False
    deck.para(s, MARGIN_L, BODY_TOP + 3.92, CONTENT_W,
              "这节课的每一页，右上角标签都会标出它落在哪一段路线上。",
              pt=20, color=MUTED, align="ctr", spacing=1.0)
    deck.notes(s, "路线图只讲 1 分钟，作用是给学生「我在哪」的坐标。")

    # ── P4 概念理解（左问题 / 右矢量图）──
    s = deck.blank("概念理解")
    deck.header(s, "先看图：把曲边切成一叠矩形", kicker="直观理解", tag="概念")
    lw = 5.55
    ix = MARGIN_L + lw + 0.44
    iw = SLIDE_W_IN - MARGIN_R - ix
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, lw, "观察右图，回答三个问题：", pt=21, color=INK, bold=True)
    y += 0.10
    qs = ["矩形越多，误差是变大还是变小？",
          "每条矩形的宽是多少？高取的是哪个点？",
          "把矩形条数推到无穷，会得到什么？"]
    for i, q in enumerate(qs):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGIN_L + 0.02), Inches(y + 0.08),
                               Inches(0.34), Inches(0.34))
        d.fill.solid(); d.fill.fore_color.rgb = CYAN
        d.line.fill.background(); d.shadow.inherit = False
        d.name = "tag_num"
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf = d.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        y += deck.para(s, MARGIN_L + 0.44, y, lw - 0.44, q, pt=20, color=INK, spacing=1.28)
        y += 0.18

    cardy, cardh = 5.05, 1.65
    deck.card(s, MARGIN_L, cardy, lw, cardh, fill=BG2, line=RULE, bar=AMBER)
    deck.label(s, MARGIN_L + 0.20, cardy + 0.14, lw - 0.44, "记住这句话", pt=16, color=AMBER)
    deck.para(s, MARGIN_L + 0.18, cardy + 0.52, lw - 0.36,
              "近似不是将就：只要让分割无限细，误差就被挤成零。",
              pt=20, color=INK, spacing=1.30)

    riemann_diagram(s, ix, BODY_TOP, iw, 2.92, n=9)
    deck.label(s, ix, BODY_TOP + 3.02, iw, "图：用 9 条矩形近似 y = x² 下方的曲边梯形",
               pt=16, color=MUTED, bold=False, align="ctr", role="cap")
    deck.card(s, ix, cardy, iw, cardh, fill=WHITE, line=RULE, bar=CYAN)
    deck.para(s, ix + 0.18, cardy + 0.18, iw - 0.36,
              "矩形总面积 = 各条「高 × 宽」之和。条数越多越贴合曲线，多出的小三角越薄。",
              pt=20, color=INK_SOFT, spacing=1.30)
    deck.notes(s, "对着图数矩形：宽都是 Δx，高是 f(ξᵢ)。把「和式」三个字写到黑板上。约 4 分钟。")

    # ── P5 核心公式 ──
    s = deck.blank("核心公式")
    deck.header(s, "把「切—乘—加—逼近」写成一行", kicker="公式表达", tag="核心公式",
                tag_color=CORAL)
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W, "四步法的最终产物，就是定积分的定义：", pt=21)
    y += 0.10
    deck.card(s, MARGIN_L, y, CONTENT_W, 1.62, fill=BG2, line=CORAL, bar=CORAL)
    pool.place(s, T["define"], MARGIN_L + 0.2, y, max_w=CONTENT_W - 0.4, max_h=1.62,
               align="ctr", valign="ctr", name="tex_define")
    y += 1.62 + 0.26
    y += deck.para(s, MARGIN_L, y, CONTENT_W, "这一行里藏着三个必须逐字读懂的部件：",
                   pt=21, color=INK, bold=True)
    y += 0.12
    parts = [("分割", "每条矩形的宽", CYAN), ("近似", "高 × 宽 = 一条的面积", AMBER),
             ("取极限", "求和之后再逼近", CORAL)]
    cw2 = (CONTENT_W - 0.5) / 3
    hcard = BODY_BOTTOM - y
    for i, (a, b, c) in enumerate(parts):
        cx = MARGIN_L + i * (cw2 + 0.25)
        deck.card(s, cx, y, cw2, hcard, fill=WHITE, line=RULE, bar=c)
        deck.label(s, cx + 0.18, y + 0.14, cw2 - 0.36, a, pt=17, color=c)
        cap_h = deck.measure_text_h(b, 20, cw2 - 0.32, spacing=1.0, space_after=0)
        pool.place(s, T[["k1", "k2", "k3"][i]], cx + 0.12, y + 0.52,
                   max_w=cw2 - 0.24, max_h=hcard - 0.62 - cap_h - 0.14,
                   align="ctr", valign="ctr", min_pt=17, name=f"tex_k{i}")
        deck.para(s, cx + 0.16, y + hcard - cap_h - 0.12, cw2 - 0.32, b, pt=20,
                  color=INK_SOFT, align="ctr", spacing=1.0)
    deck.notes(s, "逐字念定义：先有 Σ 再有 lim，顺序不能倒。强调 dx 不是装饰，"
                  "它就是那个「宽」。约 5 分钟。")

    # ── P6 概念拆解 ──
    s = deck.blank("概念拆解")
    deck.header(s, "三个关键词，读懂整条定义", kicker="公式表达", tag="拆解")
    kws = [
        ("和式", "Σ 是本体",
         "定积分首先是一个和：把 n 条矩形的面积加起来，得到一个只依赖 n 的数 Sₙ。"),
        ("极限", "lim 才是灵魂",
         "让 n→∞、Δx→0，Sₙ 收敛到唯一的那个数——这个极限值才叫定积分。"),
        ("任取", "ξᵢ 怎么取都行",
         "取左端点、右端点还是中点，极限都相同。正因为如此，这个数只由 f 与 [a,b] 决定。"),
    ]
    ch = (BODY_BOTTOM - BODY_TOP - 0.44) / 3
    yy = BODY_TOP
    for i, (k, sub, body) in enumerate(kws):
        col = [CYAN, AMBER, CORAL][i]
        deck.card(s, MARGIN_L, yy, CONTENT_W, ch, fill=WHITE if i % 2 else BG2, line=RULE, bar=col)
        deck.text(s, MARGIN_L + 0.24, yy, 1.9, ch, k, pt=28, color=col, bold=True,
                  spacing=1.0, space_after=0, anchor="ctr", role="h1")
        deck.para(s, MARGIN_L + 2.3, yy + 0.24, 3.0, sub, pt=21, color=INK, bold=True, spacing=1.0)
        bh2 = deck.measure_text_h(body, 20, CONTENT_W - 5.6, spacing=1.30, space_after=0)
        deck.para(s, MARGIN_L + 5.3, yy + (ch - bh2) / 2, CONTENT_W - 5.55, body,
                  pt=20, color=INK_SOFT, spacing=1.30)
        yy += ch + 0.22
    deck.notes(s, "三个关键词各一句。第三条最容易被忽略，但它是「定积分是个确定的数」的根据。")

    # ── P7~P16 两道逐步例题 ──
    probs[0].emit(deck)
    probs[1].emit(deck)

    # ── 核心定理：N-L 公式 ──
    s = deck.blank("核心定理")
    deck.header(s, "基本定理：把极限换成一次减法", kicker="公式表达", tag="定理", tag_color=CORAL)
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "把上限放开造出面积函数 A(x)，对它求导，就得到微分与积分互逆：", pt=21)
    y += 0.10
    deck.card(s, MARGIN_L, y, CONTENT_W, 1.20, fill=WHITE, line=RULE, bar=CYAN)
    pool.place(s, T["area_fn"], MARGIN_L + 0.2, y, max_w=CONTENT_W - 0.4, max_h=1.20,
               align="ctr", valign="ctr", name="tex_area_fn")
    y += 1.20 + 0.24
    deck.card(s, MARGIN_L, y, CONTENT_W, 1.55, fill=BG2, line=CORAL, bar=CORAL)
    deck.label(s, MARGIN_L + 0.20, y + 0.12, 5.0, "牛顿-莱布尼茨公式", pt=17, color=CORAL)
    pool.place(s, T["nl"], MARGIN_L + 0.2, y + 0.46, max_w=CONTENT_W - 0.4, max_h=1.00,
               align="ctr", valign="ctr", name="tex_nl")
    y += 1.55 + 0.24
    deck.para(s, MARGIN_L, y, CONTENT_W,
              "对照例 1：用定义算了整整五步，用这条公式只需 [x³/3]₀¹ = 1/3 一行。"
              "这就是它被称为「基本定理」的原因。",
              pt=20, color=INK_SOFT, spacing=1.30)
    deck.notes(s, "务必回扣例 1：同一个积分，两种算法，一难一易。让学生自己感叹一下。约 5 分钟。")

    # ── 方法清单 ──
    s = deck.blank("方法清单")
    deck.header(s, "用 N-L 公式算定积分的四步", kicker="方法提炼", tag="方法")
    steps4 = [
        ("第一步｜判号拆区间", "被积函数带绝对值或分段时，先找分界点，按正负拆开。"),
        ("第二步｜找原函数", "对每段找 F(x) 使 F′(x)=f(x)，求导验回去确认没找错。"),
        ("第三步｜代上下限", "写成 [F(x)]ₐᵇ，再算 F(b)−F(a)，上限减下限不能写反。"),
        ("第四步｜回题意", "问面积就取绝对值，问位移就带符号，问路程再取绝对值。"),
    ]
    cw3 = (CONTENT_W - 0.75) / 4
    ytop = BODY_TOP + 0.16
    hh3 = max(deck.measure_text_h(h, 20, cw3 - 0.32, spacing=1.05, space_after=0, bold=True)
              for h, _b in steps4)
    bh3 = max(deck.measure_text_h(b, 20, cw3 - 0.32, spacing=1.30, space_after=0)
              for _h, b in steps4)
    body_y = 0.92 + hh3 + 0.10
    cardh = body_y + bh3 + 0.20
    for i, (h1, b1) in enumerate(steps4):
        cx = MARGIN_L + i * (cw3 + 0.25)
        deck.card(s, cx, ytop, cw3, cardh, fill=BG2, line=RULE)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.18), Inches(ytop + 0.22),
                                 Inches(0.52), Inches(0.52))
        num.fill.solid(); num.fill.fore_color.rgb = [CYAN, AMBER, CORAL, INK][i]
        num.line.fill.background(); num.shadow.inherit = False
        num.name = "tag_num"
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf = num.text_frame; tf.word_wrap = False
        p_ = tf.paragraphs[0]; p_.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p_.add_run(); r.text = str(i + 1)
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        deck.para(s, cx + 0.16, ytop + 0.92, cw3 - 0.32, h1, pt=20, color=INK, bold=True,
                  spacing=1.05)
        deck.para(s, cx + 0.16, ytop + body_y, cw3 - 0.32, b1, pt=20, color=INK_SOFT,
                  spacing=1.30)
    cy = ytop + cardh + 0.22
    deck.card(s, MARGIN_L, cy, CONTENT_W, BODY_BOTTOM - cy, fill=WHITE, line=CORAL, bar=CORAL)
    deck.para(s, MARGIN_L + 0.22, cy + 0.18, CONTENT_W - 0.44,
              "一句口诀：先判号、再拆段；找原函数、代上下限；最后一定回头看题问的是什么。",
              pt=21, color=INK, bold=True, spacing=1.0)
    deck.notes(s, "四步固定动作，作业按这四步给分。第四步是本讲最容易丢分的地方。")

    # ── 易错辨析 ──
    s = deck.blank("易错辨析")
    deck.header(s, "三个最常见的错，一次只看一个", kicker="易错辨析", tag="易错", tag_color=CORAL)
    errs = [
        ("错法一", "上下限写反还不变号", "算完 F(a)−F(b) 就直接抄答案。",
         "交换上下限必须变号；写成 [F(x)]ₐᵇ 就不会记错顺序。", "err1"),
        ("错法二", "带绝对值不拆区间", "把 |x²−1| 当成 x²−1 一路积到底。",
         "先解 f(x)=0 找到分界点，再按正负拆成两段分别积。", "err2"),
        ("错法三", "把定积分当成面积", "曲线在 x 轴下方时，定积分是负的。",
         "问面积要用 ∫|f|；问位移才用带符号的 ∫f。", "err3"),
    ]
    cwe = (CONTENT_W - 0.5) / 3
    che = BODY_BOTTOM - BODY_TOP
    for i, (tag, head, wrong, right, tk) in enumerate(errs):
        cx = MARGIN_L + i * (cwe + 0.25)
        deck.card(s, cx, BODY_TOP, cwe, che, fill=WHITE, line=RULE, bar=CORAL)
        yy = BODY_TOP + 0.12
        yy += deck.label(s, cx + 0.20, yy, cwe - 0.40, tag, pt=16, color=CORAL) + 0.04
        yy += deck.para(s, cx + 0.18, yy, cwe - 0.36, head, pt=21, color=INK, bold=True,
                        spacing=1.16) + 0.08
        yy += deck.para(s, cx + 0.18, yy, cwe - 0.36, wrong, pt=20, color=MUTED,
                        spacing=1.28) + 0.10
        pool.place(s, T[tk], cx + 0.14, yy, max_w=cwe - 0.28, max_h=0.88, align="ctr",
                   valign="ctr", min_pt=15, name=f"tex_{tk}")
        yy += 0.88 + 0.14
        deck.rule(s, cx + 0.24, yy, cwe - 0.48)
        yy += 0.16
        yy += deck.label(s, cx + 0.20, yy, 1.4, "正解", pt=16, color=CYAN) + 0.04
        deck.para(s, cx + 0.18, yy, cwe - 0.36, right, pt=20, color=CYAN, spacing=1.28)
    deck.notes(s, "一卡一错，先念错法让学生举手。错法二正是例 2 的考点，错法三是应用题的分水岭。")

    # ── 迁移应用 ──
    s = deck.blank("迁移应用")
    deck.header(s, "同一个累积结构，换个场景照用", kicker="迁移应用", tag="迁移")
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "把「无数个微小量累加起来」的思路搬到物理，就得到两个高频应用：", pt=21)
    y += 0.12
    half = (CONTENT_W - 0.40) / 2
    hh = 2.10
    deck.card(s, MARGIN_L, y, half, hh, fill=BG2, line=RULE, bar=CYAN)
    deck.label(s, MARGIN_L + 0.20, y + 0.14, half - 0.4, "变速直线运动", pt=17, color=CYAN)
    pool.place(s, T["move"], MARGIN_L + 0.14, y + 0.50, max_w=half - 0.28, max_h=0.90,
               align="ctr", valign="ctr", min_pt=16, name="tex_move")
    deck.para(s, MARGIN_L + 0.18, y + 1.48, half - 0.36,
              "位移带符号，路程取绝对值——两者不等就说明中途掉了头。",
              pt=20, color=INK_SOFT, spacing=1.26)
    rx = MARGIN_L + half + 0.40
    deck.card(s, rx, y, half, hh, fill=BG2, line=RULE, bar=CORAL)
    deck.label(s, rx + 0.20, y + 0.14, half - 0.4, "变力做功", pt=17, color=CORAL)
    pool.place(s, T["work"], rx + 0.14, y + 0.50, max_w=half - 0.28, max_h=0.90,
               align="ctr", valign="ctr", min_pt=16, name="tex_work")
    deck.para(s, rx + 0.18, y + 1.48, half - 0.36,
              "力随位置变化时，W = F·s 失效，改成把 F(x)dx 一段段累加。",
              pt=20, color=INK_SOFT, spacing=1.26)
    y += hh + 0.26
    tip = "共同点：都是「微小量 × 微小区间，再累加取极限」。认出这个结构，题目就只剩计算。"
    tip_h = deck.measure_text_h(tip, 21, CONTENT_W - 0.44, spacing=1.10, space_after=0, bold=True)
    deck.card(s, MARGIN_L, y, CONTENT_W, tip_h + 0.36, fill=WHITE, line=AMBER, bar=AMBER)
    deck.para(s, MARGIN_L + 0.22, y + 0.18, CONTENT_W - 0.44, tip, pt=21, color=INK,
              bold=True, spacing=1.10)
    deck.notes(s, "弹簧做功 W=∫kx dx 可当堂口算一道。强调「认结构」比「记公式」重要。约 4 分钟。")

    # ── 课堂练习 ──
    s = deck.blank("课堂练习")
    deck.header(s, "两道练习：先说方法，再写表达", kicker="课堂练习", tag="练习", tag_color=AMBER)
    exs = [
        ("练习 1｜基础", "用 N-L 公式计算下列定积分。", "ex1",
         "先说方法：3x²+2x 的原函数是什么？",
         "再写表达：写成 [F(x)]₀¹，代上限减下限。"),
        ("练习 2｜进阶", "质点速度如下，求 0 到 3 秒的位移与路程。", "ex2",
         "先说方法：v(t)=0 的根在哪？哪段速度为负？",
         "再写表达：位移直接积，路程分段取绝对值再相加。"),
    ]
    cw4 = (CONTENT_W - 0.4) / 2
    for i, (tag, ask, tk, m1, m2) in enumerate(exs):
        cx = MARGIN_L + i * (cw4 + 0.4)
        deck.card(s, cx, BODY_TOP, cw4, BODY_BOTTOM - BODY_TOP, fill=BG2 if i == 0 else WHITE,
                  line=RULE, bar=[CYAN, CORAL][i])
        yy = BODY_TOP + 0.18
        yy += deck.label(s, cx + 0.20, yy, cw4 - 0.4, tag, pt=17, color=[CYAN, CORAL][i]) + 0.08
        yy += deck.para(s, cx + 0.18, yy, cw4 - 0.36, ask, pt=21, color=INK, spacing=1.24) + 0.14
        boxh = 1.32
        pool.place(s, T[tk], cx + 0.16, yy, max_w=cw4 - 0.32, max_h=boxh, align="ctr",
                   valign="ctr", min_pt=17, name=f"tex_{tk}")
        yy += boxh + 0.24
        for j, m in enumerate((m1, m2)):
            deck.rule(s, cx + 0.22, yy + 0.14, 0.07, color=[CYAN, CORAL][i], thick=0.10)
            yy += deck.para(s, cx + 0.36, yy, cw4 - 0.56, m, pt=20,
                            color=INK_SOFT if j == 0 else MUTED, spacing=1.26) + 0.16
    deck.notes(s, "练习 1 全体动笔 2 分钟；练习 2 找学生上台，重点看他有没有先判号再拆段。约 8 分钟。")

    # ── 课堂总结 ──
    s = deck.blank("课堂总结")
    deck.header(s, "三句话带走这节课", kicker="课堂总结", tag="小结")
    lines3 = [
        ("其一", "定积分是和式的极限。", "先分割求和得到 Sₙ，再让 n→∞，收敛到的那个数才是它。"),
        ("其二", "基本定理是把极限换减法。", "找到原函数 F，一次 F(b)−F(a) 就替掉了整个求极限过程。"),
        ("其三", "算完一定回头看题。", "问面积取绝对值，问位移带符号——这一步比计算更容易丢分。"),
    ]
    yy = BODY_TOP
    for i, (n, a, b) in enumerate(lines3):
        col = [CYAN, AMBER, CORAL][i]
        deck.label(s, MARGIN_L, yy + 0.06, 1.05, n, pt=18, color=col)
        deck.para(s, MARGIN_L + 0.96, yy, 5.5, a, pt=22, color=INK, bold=True, spacing=1.0)
        bh4 = deck.measure_text_h(b, 20, CONTENT_W - 6.5, spacing=1.26, space_after=0)
        deck.para(s, MARGIN_L + 6.5, yy, CONTENT_W - 6.5, b, pt=20, color=INK_SOFT, spacing=1.26)
        yy += max(0.60, bh4) + 0.18
        if i < 2:
            deck.rule(s, MARGIN_L, yy - 0.10, CONTENT_W)
    cy = yy + 0.18
    deck.card(s, MARGIN_L, cy, CONTENT_W, BODY_BOTTOM - cy, fill=BG2, line=CORAL, bar=CORAL)
    deck.label(s, MARGIN_L + 0.22, cy + 0.16, 4.0, "离堂检验（交条子）", pt=17, color=CORAL)
    deck.para(s, MARGIN_L + 0.18, cy + 0.56, 8.3,
              "算出它的值，并回答：它等于曲线与 x 轴围成的面积吗？说明理由。",
              pt=20, color=INK, spacing=1.26)
    pool.place(s, T["out1"], MARGIN_L + 8.7, cy + 0.18, max_w=CONTENT_W - 8.9,
               max_h=BODY_BOTTOM - cy - 0.36, align="ctr", valign="ctr", name="tex_out1")
    deck.notes(s, "答案是 0，但面积是 1——这个反差正好回收错法三。下节课讲定积分求两曲线间面积。")


def main():
    os.makedirs(OUT, exist_ok=True)
    from genimg import gen_all
    gen_all(IMAGES, IMG)
    pool = TexPool(TEXDIR, scale=4)
    T = register(pool)
    probs = problems(pool)
    pool.render()
    deck = Deck(footer_left=FOOTER, pool=pool)
    build(deck, pool, T, probs)
    path = deck.save(os.path.join(OUT, "定积分与微积分基本定理.pptx"))
    print(f"课件C: {len(deck.pages)} 页 -> {path}")
    return path


if __name__ == "__main__":
    main()
