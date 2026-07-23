# -*- coding: utf-8 -*-
"""课件A：《函数的极限与连续性 —— 从“无限接近”到 ε-δ》

内容来源：mathkit/build_10.py（第10讲 函数的极限与连续性）。原课件只作知识来源，不被修改。
选题理由：对十讲做公式复杂度审查，第10讲难点权重最高（极限15 + 量词8 + 分式8 + 多层上下标16）。

产出：out/limit-continuity/函数的极限与连续性.pptx + 公式PNG + 概念图 + 逐页实渲预览
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
KIT = os.path.join(HERE, "..", "texkit")
sys.path.insert(0, os.path.abspath(KIT))

from pptx.util import Inches, Pt                                    # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                              # noqa: E402
from theme import (INK, INK_SOFT, MUTED, CYAN, CORAL, AMBER, BG, BG2, RULE, WHITE,  # noqa: E402
                   FONT, SLIDE_W_IN, SLIDE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
                   BODY_TOP, BODY_BOTTOM, FOOTER_Y, PT_TITLE, PT_H2, PT_BODY,
                   PT_TEX_DISPLAY, H_INK, H_CYAN, H_CORAL, H_AMBER)
from slides import Deck, LayoutError                                # noqa: E402
from tex import TexPool                                             # noqa: E402
from stepdeck import StepProblem                                    # noqa: E402
from diagram import hole_diagram, squeeze_diagram                   # noqa: E402

OUT = os.path.abspath(os.path.join(HERE, "..", "out", "limit-continuity"))
IMG = os.path.join(OUT, "img")
TEXDIR = os.path.join(OUT, "tex")
FOOTER = "高中数学 · 函数的极限与连续性 · ε-δ 与分段函数"

# ────────────────────────────── 概念图（文生图）──────────────────────────────
# 说明：承载数学关系的示意图一律走 texkit/diagram.py 矢量绘制（实测文生图画不对数学关系）。
# 文生图只负责不含数学关系的**情境意象**，本课件用在封面。
IMAGES = {
    "approach": ("deep indigo and teal abstract illustration, a long corridor receding "
                 "into the distance, a row of glowing arch frames along it spaced closer "
                 "and closer together as they approach a small bright point of light far "
                 "ahead, cinematic minimal, dark background, soft glow, "
                 "no people, no text, no letters, no numbers"),
}


# ────────────────────────────── 公式清单 ──────────────────────────────
def register(pool):
    T = {}
    A = lambda k, tex, pt=PT_TEX_DISPLAY, c=H_INK: T.__setitem__(k, pool.add(tex, pt=pt, color=c))

    # 问题驱动
    A("q_hole", r"f(x)=\frac{x^{2}-1}{x-1}\quad(x\neq 1)", 34)
    A("q_val", r"x\to 1\ \text{时}\ f(x)\to\ ?", 30, H_CORAL)
    # 直观描述
    A("intuit", r"\lim_{x\to a}f(x)=L", 40, H_INK)
    # 核心公式：ε-δ
    A("epsdelta",
      r"\lim_{x\to a}f(x)=L\iff\forall\,\varepsilon>0,\ \exists\,\delta>0,\ "
      r"\text{使}\ 0<|x-a|<\delta\Rightarrow|f(x)-L|<\varepsilon", 26, H_INK)
    A("k1", r"\forall\,\varepsilon>0", 30, H_CYAN)
    A("k2", r"\exists\,\delta>0", 30, H_CYAN)
    A("k3", r"0<|x-a|<\delta", 30, H_CYAN)
    # 左右极限 / 连续
    A("lr", r"\lim_{x\to a}f(x)=L\iff\lim_{x\to a^{-}}f(x)=\lim_{x\to a^{+}}f(x)=L", 26)
    A("cont", r"f\ \text{在}\ x=a\ \text{处连续}\iff\lim_{x\to a}f(x)=f(a)", 27)
    # 迁移
    A("squeeze",
      r"g(x)\le f(x)\le h(x)\ \text{且}\ \lim_{x\to a}g(x)=\lim_{x\to a}h(x)=L"
      r"\ \Rightarrow\ \lim_{x\to a}f(x)=L", 22)
    A("lim1", r"\lim_{x\to 0}\frac{\sin x}{x}=1", 34, H_CORAL)
    A("lim2", r"\lim_{n\to\infty}\left(1+\frac{1}{n}\right)^{\!n}=e", 30, H_CORAL)
    # 易错
    A("err1", r"\lim_{x\to1}\frac{x^{2}-1}{x-1}=2", 22, H_CORAL)
    A("err2", r"\lim_{x\to0^{-}}\frac{|x|}{x}=-1\neq 1", 22, H_CORAL)
    A("err3", r"\delta=\frac{\varepsilon}{3}\ \text{或}\ \frac{\varepsilon}{4}", 22, H_CORAL)
    # 练习
    A("ex1", r"\lim_{x\to 3}(4x-5)=7", 26)
    A("ex2", r"f(x)=\begin{cases}\dfrac{\sqrt{x+1}-1}{x}, & x>0\\[4pt] a x+2, & x\le 0\end{cases}", 24)
    # 总结/检验
    A("out1", r"\lim_{x\to 2}\frac{x^{2}-4}{x-2}", 28)
    return T


# ────────────────────────────── 两道逐步例题 ──────────────────────────────
def problems(pool):
    p1 = StepProblem(
        pool, no=1, kind="用定义证明极限",
        title="用 ε-δ 定义证明 lim(3x−1)=5",
        stem="用极限的 ε-δ 定义证明：当 x→2 时，3x−1 的极限等于 5。",
        stem_note="定义题的本质不是“算”，而是“对任意给定的精度 ε，都能交出一个 δ”。",
        steps=[
            {"head": "写出要控制的目标量", "tex": [r"|f(x)-L|=\bigl|(3x-1)-5\bigr|=|3x-6|=3|x-2|"],
             "note": "把差化成含 |x−2| 的形式。",
             "speak": "所有 ε-δ 证明的第一步，都是把 |f(x)−L| 化成 |x−a| 的倍数或可放大式；"
                      "这一步直接决定了 δ 怎么取。"},
            {"head": "由目标量倒推 δ", "tex": [r"3|x-2|<\varepsilon\iff|x-2|<\frac{\varepsilon}{3}"],
             "note": "倒推只是草稿工作。",
             "speak": "想让左边小于 ε，只要 |x−2| 小于 ε/3 就够。强调这一步写在草稿上，"
                      "正式书写时要“正着写”。"},
            {"head": "取定 δ 并正式书写", "tex": [r"\text{取}\ \delta=\frac{\varepsilon}{3}>0"],
             "note": "δ 存在一个即可，不必最大。",
             "speak": "取 δ = ε/3 或 ε/4 都对。此处回收易错点三：δ 不唯一，但必须依赖 ε。"},
            {"head": "验证：由 δ 推出 ε", "tex": [
                r"0<|x-2|<\delta\Rightarrow|f(x)-5|=3|x-2|<3\delta=\varepsilon"],
             "note": "方向必须是由 δ 推出 ε。",
             "speak": "这一行才是真正的证明。让学生对照：草稿倒着想，答卷正着写。"},
            {"head": "结论", "tex": [r"\forall\,\varepsilon>0,\ \exists\,\delta=\tfrac{\varepsilon}{3}>0,"
                                     r"\ 0<|x-2|<\delta\Rightarrow|(3x-1)-5|<\varepsilon"],
             "note": "按定义语序抄一遍，证毕。", "final": True,
             "speak": "板书示范：量词顺序 ∀ε 在前、∃δ 在后，颠倒即错。"},
        ])

    p2 = StepProblem(
        pool, no=2, kind="分段函数的连续性",
        title="分段函数在分界点连续，求参数 a",
        stem="已知 f(x) 在 x = 1 处连续，求实数 a 的值。",
        stem_note="分界点处的连续性 = 左极限 = 右极限 = 函数值，三者缺一不可。",
        steps=[
            {"head": "写清函数与待定参数", "tex": [
                r"f(x)=\begin{cases}\dfrac{x^{2}-1}{x-1}, & x\neq 1\\[4pt] a, & x=1\end{cases}"],
             "note": "x≠1 可约分，x=1 由 a 指定。",
             "speak": "先让学生指出：这是可去间断点的典型结构——函数值被单独钉了一个 a。"},
            {"head": "求 x→1 的极限（约分去掉 0/0）", "tex": [
                r"\lim_{x\to1}\frac{x^{2}-1}{x-1}=\lim_{x\to1}\frac{(x+1)(x-1)}{x-1}=\lim_{x\to1}(x+1)=2"],
             "note": "约分合法，靠的是“去心”。",
             "speak": "讨论极限时 x≠1，分母不为零，所以可以约分。"
                      "回扣核心公式里 0<|x−a| 这个去心条件——这是本页最关键的一句。"},
            {"head": "验证左右极限一致", "tex": [
                r"\lim_{x\to1^{-}}(x+1)=2=\lim_{x\to1^{+}}(x+1)"],
             "note": "两侧同一解析式，自动相等。",
             "speak": "两侧用的是同一个式子，故左右极限自动相等，极限存在。"
                      "对比反例 |x|/x：两侧解析式不同，就必须分别算。"},
            {"head": "套用连续的三条件", "tex": [
                r"f\ \text{在}\ x=1\ \text{连续}\iff\lim_{x\to1}f(x)=f(1)=a"],
             "note": "三条件里的第三条钉死 a。",
             "speak": "函数在该点有定义、极限存在、二者相等——第三条把 a 钉死。"
                      "把连续定义写在黑板右上角，随时对照。"},
            {"head": "解出参数", "tex": [r"a=2"],
             "note": "洞被填平，间断点消除。", "final": True,
             "speak": "总结：可去间断点之所以“可去”，就是因为能这样补一个值。"},
        ])
    return [p1, p2]


# ────────────────────────────── 建页 ──────────────────────────────
def build(deck, pool, T, probs):
    # ── P1 封面 ──
    s = deck.blank("封面", bg=INK, footer=False)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(2.30),
                              Inches(0.14), Inches(1.62))
    band.fill.solid(); band.fill.fore_color.rgb = AMBER
    band.line.fill.background(); band.shadow.inherit = False
    tw = 7.05
    cover = os.path.join(IMG, "approach.png")
    if os.path.exists(cover):
        deck.picture(s, cover, MARGIN_L + 0.34 + tw + 0.40, 2.16,
                     SLIDE_W_IN - MARGIN_R - (MARGIN_L + 0.34 + tw + 0.40), 2.66, mode="cover")
    deck.label(s, MARGIN_L + 0.34, 1.72, tw, "高中数学 · 强基衔接", pt=18, color=CYAN)
    deck.para(s, MARGIN_L + 0.34, 2.20, tw, "函数的极限与连续性", pt=46, color=WHITE,
              bold=True, spacing=1.05, role="h1")
    deck.para(s, MARGIN_L + 0.34, 3.34, tw, "教学主线：直观理解 → 公式表达 → 迁移应用",
              pt=22, color=RULE, spacing=1.0)
    deck.rule(s, MARGIN_L + 0.34, 4.16, 3.2, color=AMBER, thick=0.03)
    deck.para(s, MARGIN_L + 0.34, 4.42, tw,
              "一句核心理解：极限只刻画 x 趋近 a 的过程，与 f(a) 的取值毫无关系。",
              pt=24, color=WHITE, spacing=1.25)
    deck.label(s, MARGIN_L + 0.34, 6.52, 11.0, FOOTER, pt=15, color=MUTED, bold=False)
    deck.notes(s, "开场 1 分钟：先不给定义，只抛出“无限接近到底是什么意思”。"
                  "本节课的主线写在标题下方，全程回指。")

    # ── P2 问题驱动 ──
    s = deck.blank("问题驱动")
    deck.header(s, "一个填不上的洞", kicker="问题驱动", tag="导入")
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "把下面这个函数在 x = 1 附近画出来，你会看到一条完整的直线上缺了一个点：",
                   pt=21)
    y += 0.10
    pool.place(s, T["q_hole"], MARGIN_L, y, max_w=CONTENT_W, align="ctr", name="tex_q_hole")
    y += pool.size_in(T["q_hole"])[1] + 0.10
    pool.place(s, T["q_val"], MARGIN_L, y, max_w=CONTENT_W, align="ctr", name="tex_q_val")
    y += pool.size_in(T["q_val"])[1] + 0.22

    cw0 = (CONTENT_W - 0.5) / 3
    cards2 = [
        ("函数值不存在", "x = 1 代进去分母为 0，这点没有函数值。", CYAN),
        ("趋势却很确定", "x 越靠近 1，f(x) 就越靠近 2。", AMBER),
        ("于是必须分家", "“某点的值”与“附近的趋势”是两件事。", CORAL),
    ]
    bh = max(deck.measure_text_h(b, 20, cw0 - 0.40, spacing=1.26, space_after=0)
             for _h, b, _c in cards2)
    cardh = 0.62 + bh + 0.20
    for i, (h1, b1, col) in enumerate(cards2):
        cx = MARGIN_L + i * (cw0 + 0.25)
        deck.card(s, cx, y, cw0, cardh, fill=WHITE, line=RULE, bar=col)
        deck.para(s, cx + 0.20, y + 0.16, cw0 - 0.40, h1, pt=21, color=col, bold=True,
                  spacing=1.0)
        deck.para(s, cx + 0.20, y + 0.62, cw0 - 0.40, b1, pt=20, color=INK_SOFT, spacing=1.26)
    deck.notes(s, "让学生口算 x=0.9/0.99/0.999 的三个值，自己说出“越来越接近 2”。"
                  "不要急着给“极限”这个词，先让直觉到位。约 4 分钟。")

    # ── P3 学习路线 ──
    s = deck.blank("学习路线")
    deck.header(s, "本节课的三段路线", kicker="学习路线", tag="路线")
    routes = [
        ("① 直观理解", "无限接近", "从数值表和图像上的“洞”出发，先建立“趋势”的感觉，暂不追求严格。"),
        ("② 公式表达", "ε-δ 语言", "把“无限接近”翻译成不等式：任给精度 ε，都能找出范围 δ。"),
        ("③ 迁移应用", "连续与夹逼", "用极限判连续、定参数，再迁移到夹逼定理与两个重要极限。"),
    ]
    cw = (CONTENT_W - 0.7) / 3
    for i, (tag, head, body) in enumerate(routes):
        cx = MARGIN_L + i * (cw + 0.35)
        col = [CYAN, AMBER, CORAL][i]
        deck.card(s, cx, BODY_TOP + 0.24, cw, 3.36, fill=BG2 if i % 2 == 0 else WHITE,
                  line=RULE, bar=col)
        deck.label(s, cx + 0.22, BODY_TOP + 0.44, cw - 0.44, tag, pt=17, color=col)
        deck.para(s, cx + 0.22, BODY_TOP + 0.88, cw - 0.44, head, pt=28, color=INK,
                  bold=True, spacing=1.05)
        deck.para(s, cx + 0.22, BODY_TOP + 1.62, cw - 0.44, body, pt=20, color=INK_SOFT,
                  spacing=1.30)
        if i < 2:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx + cw + 0.06),
                                    Inches(BODY_TOP + 1.70), Inches(0.23), Inches(0.30))
            ar.fill.solid(); ar.fill.fore_color.rgb = RULE
            ar.line.fill.background(); ar.shadow.inherit = False
    deck.para(s, MARGIN_L, BODY_TOP + 3.92, CONTENT_W,
              "这节课的每一页，右上角标签都会标出它落在哪一段路线上。",
              pt=20, color=MUTED, align="ctr", spacing=1.0)
    deck.notes(s, "路线图页只讲 1 分钟，作用是给学生一个“我现在在哪”的坐标。"
                  "后面每页右上角的标签就对应这三段。")

    # ── P4 概念理解（左问题 / 右概念图）──
    s = deck.blank("概念理解")
    deck.header(s, "先看图：趋势与取值是两件事", kicker="直观理解", tag="概念")
    lw = 5.55
    ix = MARGIN_L + lw + 0.44
    iw = SLIDE_W_IN - MARGIN_R - ix
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, lw, "观察右图，回答三个问题：", pt=21, color=INK, bold=True)
    y += 0.10
    qs = ["曲线在那个空心点处断了吗？",
          "只从左边走，值趋向多少？右边呢？",
          "把空心点补实，要补多少？补错会怎样？"]
    for i, q in enumerate(qs):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGIN_L + 0.02), Inches(y + 0.08),
                               Inches(0.34), Inches(0.34))
        d.fill.solid(); d.fill.fore_color.rgb = CYAN
        d.line.fill.background(); d.shadow.inherit = False
        d.name = "tag_num"
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf = d.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        y += deck.para(s, MARGIN_L + 0.44, y, lw - 0.44, q, pt=20, color=INK, spacing=1.28)
        y += 0.18

    cardy, cardh = 5.05, 1.65
    deck.card(s, MARGIN_L, cardy, lw, cardh, fill=BG2, line=RULE, bar=AMBER)
    deck.label(s, MARGIN_L + 0.20, cardy + 0.14, lw - 0.44, "记住这句话", pt=16, color=AMBER)
    deck.para(s, MARGIN_L + 0.18, cardy + 0.52, lw - 0.36,
              "极限描述“走过去的过程”，函数值描述“站在那一点”，两者可以不一致。",
              pt=20, color=INK, spacing=1.30)

    hole_diagram(s, ix, BODY_TOP, iw, 2.92)
    deck.label(s, ix, BODY_TOP + 3.02, iw, "图：直线在一点处被挖掉，留下一个空心点",
               pt=16, color=MUTED, bold=False, align="ctr", role="cap")
    deck.card(s, ix, cardy, iw, cardh, fill=WHITE, line=RULE, bar=CYAN)
    deck.para(s, ix + 0.18, cardy + 0.18, iw - 0.36,
              "从左右两侧走向那个洞，高度都逼近同一个数——这个“共同去处”就是极限；"
              "洞本身补不补，是另一回事。",
              pt=20, color=INK_SOFT, spacing=1.30)
    deck.notes(s, "配合这幅图提三个问题，让学生用自己的话说出“趋势 ≠ 取值”。"
                  "第三问为后面的“可去间断点补值”埋伏笔。约 4 分钟。")

    # ── P5 核心公式 ──
    s = deck.blank("核心公式")
    deck.header(s, "把“无限接近”写成不等式", kicker="公式表达", tag="核心公式", tag_color=CORAL)
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "柯西与魏尔斯特拉斯给出的严格定义，就是本节课的核心公式：", pt=21)
    y += 0.10
    deck.card(s, MARGIN_L, y, CONTENT_W, 1.62, fill=BG2, line=CORAL, bar=CORAL)
    pool.place(s, T["epsdelta"], MARGIN_L + 0.2, y, max_w=CONTENT_W - 0.4, max_h=1.62,
               align="ctr", valign="ctr", name="tex_epsdelta")
    y += 1.62 + 0.26
    y += deck.para(s, MARGIN_L, y, CONTENT_W, "这一行里藏着三个必须逐字读懂的部件：",
                   pt=21, color=INK, bold=True)
    y += 0.12
    parts = [("挑战", "任意给定的精度", CYAN), ("应战", "总能找到的范围", AMBER),
             ("去心", "只看邻域不看点", CORAL)]
    cw2 = (CONTENT_W - 0.5) / 3
    hcard = BODY_BOTTOM - y
    for i, (a, b, c) in enumerate(parts):
        cx = MARGIN_L + i * (cw2 + 0.25)
        deck.card(s, cx, y, cw2, hcard, fill=WHITE, line=RULE, bar=c)
        deck.label(s, cx + 0.18, y + 0.14, cw2 - 0.36, a, pt=17, color=c)
        pool.place(s, T[["k1", "k2", "k3"][i]], cx + 0.12, y + 0.52, max_w=cw2 - 0.24,
                   max_h=0.72, align="ctr", valign="ctr", min_pt=17, name=f"tex_k{i}")
        deck.para(s, cx + 0.16, y + 1.34, cw2 - 0.32, b, pt=20, color=INK_SOFT,
                  align="ctr", spacing=1.0)
    deck.notes(s, "逐字念一遍定义，念到 ∀ε 停一下问“谁先动手”。"
                  "把 ε-δ 讲成挑战—应战的博弈：学生出 ε，老师交 δ。约 5 分钟。")

    # ── P6 概念拆解 ──
    s = deck.blank("概念拆解")
    deck.header(s, "三个关键词，读懂整条定义", kicker="公式表达", tag="拆解")
    kws = [
        ("任意", "ε 是对方出的题",
         "ε 要多小有多小，且必须对所有 ε 都成立。只对某一个 ε 成立，不叫极限。"),
        ("存在", "δ 是我方交的卷",
         "能交出一个可行的 δ 就算赢，不要求最大、不要求唯一，δ 通常写成含 ε 的式子。"),
        ("去心", "圆心那一点不参与",
         "条件写的是 0 < |x−a|，把 x = a 本身排除在外。这正是可以约分、可以不管 f(a) 的依据。"),
    ]
    ch = (BODY_BOTTOM - BODY_TOP - 0.44) / 3
    yy = BODY_TOP
    for i, (k, sub, body) in enumerate(kws):
        col = [CYAN, AMBER, CORAL][i]
        deck.card(s, MARGIN_L, yy, CONTENT_W, ch, fill=WHITE if i % 2 else BG2, line=RULE, bar=col)
        deck.text(s, MARGIN_L + 0.24, yy, 1.6, ch, k, pt=30, color=col, bold=True,
                  spacing=1.0, space_after=0, anchor="ctr", role="h1")
        deck.para(s, MARGIN_L + 2.0, yy + 0.24, 3.1, sub, pt=21, color=INK, bold=True,
                  spacing=1.0)
        bh = deck.measure_text_h(body, 20, CONTENT_W - 5.4, spacing=1.30, space_after=0)
        deck.para(s, MARGIN_L + 5.1, yy + (ch - bh) / 2, CONTENT_W - 5.35, body,
                  pt=20, color=INK_SOFT, spacing=1.30)
        yy += ch + 0.22
    deck.notes(s, "三个关键词各一句话，配合手势：ε 用手指比“很小”，δ 用手掌比“一段范围”。"
                  "去心这一条要重讲，它是后面约分合法性的全部依据。约 3 分钟。")

    # ── P7~P16 两道逐步例题 ──
    probs[0].emit(deck)
    probs[1].emit(deck)

    # ── 方法清单 ──
    s = deck.blank("方法清单")
    deck.header(s, "ε-δ 证明的四步固定动作", kicker="方法提炼", tag="方法")
    steps4 = [
        ("第一步｜化差", "把 |f(x) − L| 恒等变形成 k·|x − a|，或放大成它的形式。"),
        ("第二步｜倒推", "在草稿上解 k·|x − a| < ε，得到 |x − a| < ε/k。"),
        ("第三步｜取 δ", "令 δ = ε/k（或更小的正数），写明 δ > 0 且只依赖 ε。"),
        ("第四步｜正写", "从 0 < |x − a| < δ 推出 |f(x) − L| < ε，方向不可颠倒。"),
    ]
    cw3 = (CONTENT_W - 0.75) / 4
    ytop = BODY_TOP + 0.16
    for i, (h1, b1) in enumerate(steps4):
        cx = MARGIN_L + i * (cw3 + 0.25)
        deck.card(s, cx, ytop, cw3, 3.42, fill=BG2, line=RULE)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.18), Inches(ytop + 0.22),
                                 Inches(0.52), Inches(0.52))
        num.fill.solid(); num.fill.fore_color.rgb = [CYAN, AMBER, CORAL, INK][i]
        num.line.fill.background(); num.shadow.inherit = False
        num.name = "tag_num"
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf = num.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        deck.para(s, cx + 0.16, ytop + 0.92, cw3 - 0.32, h1, pt=21, color=INK, bold=True,
                  spacing=1.0)
        deck.para(s, cx + 0.16, ytop + 1.40, cw3 - 0.32, b1, pt=20, color=INK_SOFT,
                  spacing=1.30)
    cy = ytop + 3.42 + 0.22
    deck.card(s, MARGIN_L, cy, CONTENT_W, BODY_BOTTOM - cy, fill=WHITE, line=CORAL, bar=CORAL)
    deck.para(s, MARGIN_L + 0.22, cy + 0.20, CONTENT_W - 0.44,
              "一句口诀：草稿倒着想，答卷正着写；δ 只需一个，ε 必须任意。",
              pt=22, color=INK, bold=True, spacing=1.0)
    deck.notes(s, "四步做成固定动作，要求学生在草稿纸上按这四行分栏写。"
                  "下节课的作业批改就按这四步给分。约 3 分钟。")

    # ── 易错辨析（一卡一错）──
    s = deck.blank("易错辨析")
    deck.header(s, "三个最常见的错，一次只看一个", kicker="易错辨析", tag="易错", tag_color=CORAL)
    errs = [
        ("错法一", "把 lim 当成代入求值", "分母为零就说“算不出来”，直接卡死。",
         "极限只看去心邻域：先约分化简，代入是最后一步。", "err1"),
        ("错法二", "不查左右极限就下结论", "见函数在该点有定义，就断言极限存在。",
         "凡分界点、凡带绝对值，两侧分别算，相等才算存在。", "err2"),
        ("错法三", "非要求出最大的 δ", "花力气找最优 δ，或写出与 ε 无关的 δ。",
         "δ 存在一个即可，且必须随 ε 而定；取更小永远不错。", "err3"),
    ]
    cwe = (CONTENT_W - 0.5) / 3
    che = BODY_BOTTOM - BODY_TOP
    for i, (tag, head, wrong, right, tk) in enumerate(errs):
        cx = MARGIN_L + i * (cwe + 0.25)
        deck.card(s, cx, BODY_TOP, cwe, che, fill=WHITE, line=RULE, bar=CORAL)
        yy = BODY_TOP + 0.12
        yy += deck.label(s, cx + 0.20, yy, cwe - 0.40, tag, pt=16, color=CORAL) + 0.04
        yy += deck.para(s, cx + 0.18, yy, cwe - 0.36, head, pt=21, color=INK, bold=True,
                        spacing=1.16) + 0.08
        yy += deck.para(s, cx + 0.18, yy, cwe - 0.36, wrong, pt=20, color=MUTED,
                        spacing=1.28) + 0.10
        pool.place(s, T[tk], cx + 0.14, yy, max_w=cwe - 0.28, max_h=0.88, align="ctr",
                   valign="ctr", min_pt=15, name=f"tex_{tk}")
        yy += 0.88 + 0.14
        deck.rule(s, cx + 0.24, yy, cwe - 0.48)
        yy += 0.16
        yy += deck.label(s, cx + 0.20, yy, 1.4, "正解", pt=16, color=CYAN) + 0.04
        deck.para(s, cx + 0.18, yy, cwe - 0.36, right, pt=20, color=CYAN, spacing=1.28)
    deck.notes(s, "每张卡只讲一个错：先念错法，让学生举手“谁这样做过”，再给纠正。"
                  "错法一对应例题 2 的约分，错法三对应例题 1 的取 δ。约 4 分钟。")

    # ── 迁移应用 ──
    s = deck.blank("迁移应用")
    deck.header(s, "同一套语言，换个场景照用", kicker="迁移应用", tag="迁移")
    y = BODY_TOP
    y += deck.para(s, MARGIN_L, y, CONTENT_W,
                   "把“两边逼近同一个数”的思路放大，就得到高中最常用的两个极限工具：", pt=21)
    y += 0.10
    deck.card(s, MARGIN_L, y, CONTENT_W, 1.46, fill=BG2, line=RULE, bar=CYAN)
    deck.label(s, MARGIN_L + 0.20, y + 0.12, 3.4, "夹逼（三明治）定理", pt=17, color=CYAN)
    pool.place(s, T["squeeze"], MARGIN_L + 0.2, y + 0.48, max_w=CONTENT_W - 0.4,
               max_h=0.90, align="ctr", valign="ctr", name="tex_squeeze")
    y += 1.46 + 0.24
    ihalf = (CONTENT_W - 0.40) / 2
    imgh = BODY_BOTTOM - y - 0.42
    squeeze_diagram(s, MARGIN_L, y, ihalf, imgh)
    deck.label(s, MARGIN_L, y + imgh + 0.08, ihalf, "图：上下两条曲线把中间那条夹向同一点",
               pt=16, color=MUTED, bold=False, align="ctr", role="cap")
    rx = MARGIN_L + ihalf + 0.40
    hh = (imgh - 0.16) / 2
    deck.card(s, rx, y, ihalf, hh, fill=WHITE, line=RULE, bar=CORAL)
    pool.place(s, T["lim1"], rx + 0.12, y, max_w=ihalf - 0.24, max_h=hh, align="ctr",
               valign="ctr", name="tex_lim1")
    deck.card(s, rx, y + hh + 0.16, ihalf, hh, fill=WHITE, line=RULE, bar=CORAL)
    pool.place(s, T["lim2"], rx + 0.12, y + hh + 0.16, max_w=ihalf - 0.24, max_h=hh,
               align="ctr", valign="ctr", name="tex_lim2")
    deck.label(s, rx, y + imgh + 0.08, ihalf, "两个重要极限：都靠“夹”得到，导数里天天用",
               pt=16, color=MUTED, bold=False, align="ctr", role="cap")
    deck.notes(s, "夹逼定理只讲结构不证：三条曲线、两头同去处、中间被逼就范。"
                  "sin x / x 的几何证明放在下一课时。约 3 分钟。")

    # ── 课堂练习 ──
    s = deck.blank("课堂练习")
    deck.header(s, "两道练习：先说方法，再写表达", kicker="课堂练习", tag="练习", tag_color=AMBER)
    exs = [
        ("练习 1｜基础", "用 ε-δ 定义证明下式成立。", "ex1",
         "先说方法：|f(x) − 7| 是 |x − 3| 的几倍",
         "再写表达：按四步写完整，注意最后一步方向。"),
        ("练习 2｜进阶", "若 f(x) 在 x = 0 处连续，求实数 a。", "ex2",
         "先说方法：x → 0⁺ 用哪个式子？分子怎么有理化？",
         "再写表达：右极限、左极限与 f(0) 三者相等解 a。"),
    ]
    cw4 = (CONTENT_W - 0.4) / 2
    for i, (tag, ask, tk, m1, m2) in enumerate(exs):
        cx = MARGIN_L + i * (cw4 + 0.4)
        deck.card(s, cx, BODY_TOP, cw4, BODY_BOTTOM - BODY_TOP, fill=BG2 if i == 0 else WHITE,
                  line=RULE, bar=[CYAN, CORAL][i])
        yy = BODY_TOP + 0.18
        yy += deck.label(s, cx + 0.20, yy, cw4 - 0.4, tag, pt=17, color=[CYAN, CORAL][i]) + 0.08
        yy += deck.para(s, cx + 0.18, yy, cw4 - 0.36, ask, pt=21, color=INK, spacing=1.24) + 0.14
        boxh = 1.32
        pool.place(s, T[tk], cx + 0.16, yy, max_w=cw4 - 0.32, max_h=boxh, align="ctr",
                   valign="ctr", min_pt=17, name=f"tex_{tk}")
        yy += boxh + 0.24
        for j, m in enumerate((m1, m2)):
            deck.rule(s, cx + 0.22, yy + 0.14, 0.07, color=[CYAN, CORAL][i], thick=0.10)
            yy += deck.para(s, cx + 0.36, yy, cw4 - 0.56, m, pt=20,
                            color=INK_SOFT if j == 0 else MUTED, spacing=1.26) + 0.16
    deck.notes(s, "练习 1 全体动笔 3 分钟，找一名学生上台按四步板演。"
                  "练习 2 分子有理化是难点，提示“乘共轭”。合计约 8 分钟。")

    # ── 课堂总结 ──
    s = deck.blank("课堂总结")
    deck.header(s, "三句话带走这节课", kicker="课堂总结", tag="小结")
    lines3 = [
        ("其一", "极限是过程，不是取值。", "0 < |x − a| 这个去心条件，把 f(a) 彻底排除在讨论之外。"),
        ("其二", "定义是一场博弈。", "对方任给 ε，我方交出依赖 ε 的 δ，先后顺序不能颠倒。"),
        ("其三", "连续是极限的等号。", "极限值恰好等于函数值才叫连续，分段函数定参就是在配这个等号。"),
    ]
    yy = BODY_TOP
    for i, (n, a, b) in enumerate(lines3):
        col = [CYAN, AMBER, CORAL][i]
        deck.label(s, MARGIN_L, yy + 0.06, 0.9, n, pt=18, color=col)
        deck.para(s, MARGIN_L + 0.86, yy, 5.3, a, pt=23, color=INK, bold=True, spacing=1.0)
        bh = deck.measure_text_h(b, 20, CONTENT_W - 6.26, spacing=1.26, space_after=0)
        deck.para(s, MARGIN_L + 6.26, yy, CONTENT_W - 6.26, b, pt=20, color=INK_SOFT,
                  spacing=1.26)
        yy += max(0.60, bh) + 0.18
        if i < 2:
            deck.rule(s, MARGIN_L, yy - 0.10, CONTENT_W)
    cy = yy + 0.18
    deck.card(s, MARGIN_L, cy, CONTENT_W, BODY_BOTTOM - cy, fill=BG2, line=CORAL, bar=CORAL)
    deck.label(s, MARGIN_L + 0.22, cy + 0.16, 4.0, "离堂检验（交条子）", pt=17, color=CORAL)
    deck.para(s, MARGIN_L + 0.18, cy + 0.56, 8.3,
              "写下这个极限的值，并用一句话说明为什么可以约分——"
              "答得出理由（去心邻域内 x ≠ 2）的同学，才算真懂了今天的定义。",
              pt=20, color=INK, spacing=1.26)
    pool.place(s, T["out1"], MARGIN_L + 8.7, cy + 0.18, max_w=CONTENT_W - 8.9,
               max_h=BODY_BOTTOM - cy - 0.36, align="ctr", valign="ctr", name="tex_out1")
    deck.notes(s, "三句话让学生齐读一遍。离堂条子当场收，下节课用它分层讲评。"
                  "预告：下节课用今天的语言去证 sin x / x → 1。约 3 分钟。")


def main():
    os.makedirs(OUT, exist_ok=True)
    from genimg import gen_all
    gen_all({k: v for k, v in IMAGES.items()}, IMG)
    pool = TexPool(TEXDIR, scale=4)
    T = register(pool)
    probs = problems(pool)
    pool.render()
    deck = Deck(footer_left=FOOTER, pool=pool)
    build(deck, pool, T, probs)
    path = deck.save(os.path.join(OUT, "函数的极限与连续性.pptx"))
    print(f"课件A: {len(deck.pages)} 页 -> {path}")
    return path


if __name__ == "__main__":
    main()
