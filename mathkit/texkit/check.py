# -*- coding: utf-8 -*-
"""导出后的 PPTX 静态体检（goal 第六节 1./3.）。

重新打开**实际导出的 .pptx**（不是构建脚本内部状态）逐页逐形状检查：
  1. 越界   —— 任何形状不得越出 16:9 画布
  2. 溢出   —— 文本按真实字宽重量一遍，行数×行高不得超过文本框
  3. 贴边   —— 文本框内边距为 0 视为文字贴框
  4. 压页脚 —— 正文形状不得进入页脚安全区
  5. 变形   —— 图片显示宽高比必须等于原图（含 crop 后）宽高比
  6. 重叠   —— 文本框与图片、文本框之间的实质性重叠
  7. 字号   —— 正文字号不得低于投影下限

用法：python check.py <a.pptx> [b.pptx ...]
"""
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from slides import wrap, block_h_in                     # noqa: E402
from theme import SLIDE_W_IN, SLIDE_H_IN, FOOTER_Y, PT_MIN_BODY  # noqa: E402

EMU_IN = 914400.0
TOL = 0.03            # 英寸容差
OVERLAP_TOL = 0.06    # 重叠面积占比阈值

# 字号分层：正文必须够投影，标签/图注/页码允许小一档（但也有下限）
PT_FLOOR = {"body": PT_MIN_BODY, "h1": 24, "tag": 15, "cap": 15, "footer": 11}
# 公式 PNG 四周有透明留白（KaTeX 容器 padding 10px/14px），做重叠判定时先缩掉，
# 否则紧贴公式的说明文字会被误判成「文字压图」。
TEX_INSET_X, TEX_INSET_Y = 0.17, 0.09


def _in(v):
    return (v or 0) / EMU_IN


def _rect(sh):
    return (_in(sh.left), _in(sh.top), _in(sh.left) + _in(sh.width), _in(sh.top) + _in(sh.height))


def _overlap(a, b):
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    if w <= 0 or h <= 0:
        return 0.0
    inter = w * h
    small = min((a[2] - a[0]) * (a[3] - a[1]), (b[2] - b[0]) * (b[3] - b[1]))
    return inter / small if small > 0 else 0.0


def check_pptx(path):
    prs = Presentation(path)
    issues = []
    cw, ch = _in(prs.slide_width), _in(prs.slide_height)
    if abs(cw / ch - 16 / 9) > 0.001:
        issues.append((0, "画布", f"不是 16:9（{cw:.2f}×{ch:.2f}in）"))

    for idx, slide in enumerate(prs.slides, 1):
        texts, pics = [], []
        for sh in slide.shapes:
            name = sh.name or ""
            x0, y0, x1, y1 = _rect(sh)
            is_bg = name.startswith("bg_")
            is_footer = name.startswith("footer_")

            # 1 越界
            if not is_bg and (x0 < -TOL or y0 < -TOL or x1 > cw + TOL or y1 > ch + TOL):
                issues.append((idx, name or sh.shape_type,
                               f"越出画布：({x0:.2f},{y0:.2f})-({x1:.2f},{y1:.2f})"))

            # 图片
            if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
                pics.append((name, (x0, y0, x1, y1)))
                try:
                    iw, ih = sh.image.size
                    cl, cr = sh.crop_left or 0, sh.crop_right or 0
                    ct, cb = sh.crop_top or 0, sh.crop_bottom or 0
                    src_ar = (iw * (1 - cl - cr)) / (ih * (1 - ct - cb))
                    disp_ar = (x1 - x0) / (y1 - y0)
                    if abs(disp_ar / src_ar - 1) > 0.012:
                        issues.append((idx, name or "图片",
                                       f"变形：显示比 {disp_ar:.3f} vs 原图比 {src_ar:.3f}"))
                except Exception as e:
                    issues.append((idx, name or "图片", f"无法校验比例: {e}"))
                if not is_footer and y1 > FOOTER_Y + TOL:
                    issues.append((idx, name or "图片", f"进入页脚安全区（底 {y1:.2f} > {FOOTER_Y}）"))
                continue

            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            raw = "\n".join(p.text for p in tf.paragraphs)
            if not raw.strip():
                continue
            texts.append((name, (x0, y0, x1, y1), raw))

            # 3 贴边
            ml, mr = _in(tf.margin_left), _in(tf.margin_right)
            if ml < 0.02 or mr < 0.02:
                issues.append((idx, name or "文本", f"内边距过小（左 {ml:.3f} 右 {mr:.3f}in），文字贴框"))

            # 7 字号 + 2 溢出
            sizes, spacing = [], 1.0
            for p in tf.paragraphs:
                spacing = p.line_spacing or spacing
                for r in p.runs:
                    if r.font.size:
                        sizes.append(r.font.size.pt)
            pt = max(sizes) if sizes else 18
            bold = any(r.font.bold for p in tf.paragraphs for r in p.runs)
            role = (name.split("_")[0] if "_" in name else "body")
            floor = PT_FLOOR.get(role, PT_MIN_BODY)
            if pt < floor - 0.01:
                issues.append((idx, name or "文本",
                               f"{role} 字号 {pt:.0f}pt < {floor}pt 下限：«{raw[:24]}»"))

            inner_w = (x1 - x0) - ml - mr
            lines = []
            for p in tf.paragraphs:
                t = p.text
                lines += wrap(t, pt, inner_w, bold) if t else [""]
            sp_after = max((p.space_after.pt if p.space_after else 0) for p in tf.paragraphs) or 0
            need = block_h_in(lines, pt, spacing if spacing else 1.0, sp_after) \
                + _in(tf.margin_top) + _in(tf.margin_bottom)
            if need > (y1 - y0) + 0.05:
                issues.append((idx, name or "文本",
                               f"文字溢出框：需 {need:.2f}in > 框高 {(y1-y0):.2f}in «{raw[:24]}»"))

            # 4 压页脚
            if not is_footer and y0 + min(need, y1 - y0) > FOOTER_Y + TOL:
                issues.append((idx, name or "文本", f"正文进入页脚安全区（底 {y1:.2f} > {FOOTER_Y}）"))

        # 6 重叠
        for i, (n1, r1, _t) in enumerate(texts):
            for n2, r2 in pics:
                if str(n2).startswith("tex_"):
                    r2 = (r2[0] + TEX_INSET_X, r2[1] + TEX_INSET_Y,
                          r2[2] - TEX_INSET_X, r2[3] - TEX_INSET_Y)
                ov = _overlap(r1, r2)
                if ov > OVERLAP_TOL:
                    issues.append((idx, f"{n1 or '文本'}×{n2 or '图片'}", f"文字压图 {ov*100:.0f}%"))
            for n2, r2, _t2 in texts[i + 1:]:
                ov = _overlap(r1, r2)
                if ov > 0.12:
                    issues.append((idx, f"{n1 or '文本'}×{n2 or '文本'}", f"文本框相互重叠 {ov*100:.0f}%"))
    return issues


def main():
    if len(sys.argv) < 2:
        print("用法: python check.py <a.pptx> ...")
        sys.exit(2)
    total = 0
    for p in sys.argv[1:]:
        issues = check_pptx(p)
        total += len(issues)
        print(f"\n=== {os.path.basename(p)} ===")
        if not issues:
            print("  ✓ 全部通过（越界/溢出/贴边/压页脚/变形/重叠/字号）")
        for pg, who, msg in issues:
            print(f"  ✗ P{pg:02d} [{who}] {msg}")
    sys.exit(1 if total else 0)


if __name__ == "__main__":
    main()
