#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Polaris Forge · Python 出片桥(路线 B 的「无限版式」上层梯队)。

**为什么存在**:Rust 原生引擎(pptx_native.rs)确定性、零安装,但版式是固定那几个 +
freeform。当你想复用 build/engine.py 里已经调好的精致排版、或用 python-pptx 的完整能力
临时造任意版式时,走这里。契约与 Rust 引擎**同一份 spec JSON**,所以上层 UI/模型无需改口径。

用法(Rust forge::spec_to_pptx_sync 在 spec 顶层带 engine:"python"/"auto" 时自动调用):
    python3 pptx_bridge.py <spec.json> <out.pptx>
成功向 stdout 打印一行 JSON:{"ok":true,"slides":N,"images":M,"warnings":[...]}。
失败退非零码 + stderr 写原因(Rust 侧据此在 auto 模式回退原生引擎)。

依赖:仅 python-pptx(pip install python-pptx)。Rust 侧已先探测其存在才会调本脚本。

**扩展点**:想加版式,就在 render_slide() 的 dispatch 里加一个分支;想复用 engine.py 的
THEMES / 版式方法,把 build/ 加进 sys.path 后 import engine 即可(见文件尾 __main__ 注释)。
"""
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except Exception as e:  # noqa: BLE001
    sys.stderr.write(f"python-pptx 不可用: {e}\n")
    sys.exit(3)

# 逻辑画布 1280×720 px，1px = 9525 EMU（与 Rust 引擎同一坐标系，spec 数值可直接互通）。
PX = 9525
CANVAS_W, CANVAS_H = 1280, 720

# 内置色板：与 pptx_native.rs 的 PALETTES 气质对齐（背景两端/正文/弱化/强调/卡片/卡片描边）。
PALETTES = {
    "minimal-white": ("FFFFFF", "F2F4F7", "1A1A1A", "6B7280", "2563EB", "FFFFFF", "E5E7EB"),
    "ink-gold": ("14161B", "1E222B", "F5F5F5", "9AA0A6", "D4B06A", "232833", "3A4150"),
    "warm-paper": ("FBF7EF", "F1E9DA", "2A2420", "8A7E6E", "B5651D", "FFFFFF", "E7DcC8"),
    "forest": ("F3F7F3", "E3EFE3", "1C2B1C", "5E7A5E", "2E7D32", "FFFFFF", "D2E3D2"),
    "slate-night": ("0F1620", "18202B", "ECEFF3", "8B95A3", "38BDF8", "1B2532", "2C3A4B"),
}


def palette(name):
    return PALETTES.get(name or "", PALETTES["minimal-white"])


def parse_color(raw, pal, fallback):
    """#RRGGBB / #RGB / 裸 hex / 色板词 → RGBColor；认不出退 fallback。"""
    bg1, bg2, ink, muted, accent, card, line = pal
    tokens = {"ink": ink, "text": ink, "muted": muted, "accent": accent, "card": card,
              "line": line, "bg": bg1, "bg1": bg1, "bg2": bg2, "white": "FFFFFF", "black": "000000"}
    t = (raw or "").strip()
    if not t:
        return RGBColor.from_string(fallback)
    if t.lower() in tokens:
        return RGBColor.from_string(tokens[t.lower()])
    h = t.lstrip("#")
    if len(h) == 3 and all(c in "0123456789abcdefABCDEF" for c in h):
        h = "".join(c * 2 for c in h)
    if len(h) == 6 and all(c in "0123456789abcdefABCDEF" for c in h):
        return RGBColor.from_string(h.upper())
    return RGBColor.from_string(fallback)


def emu(px):
    return Emu(int(round(px * PX)))


def set_bg(slide, pal):
    """整页纯色背景（取渐变浅端；python-pptx 原生渐变支持弱，用实色更稳）。"""
    bg1 = pal[0]
    el = slide.background.element
    # 直接写 solidFill 到 cSld/bg（python-pptx 无高层 API）。
    from pptx.oxml import parse_xml
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    bgpr = parse_xml(
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f'<p:bgPr {ns}><a:solidFill><a:srgbClr val="{bg1}"/></a:solidFill>'
        f'<a:effectLst/></p:bgPr></p:bg>'
    )
    csld = slide.shapes._spTree.getparent()
    csld.insert(0, bgpr)


def add_text(slide, x, y, w, h, runs, align="l", anchor="t"):
    """runs: [(text, size_pt, RGBColor, bold, italic)]，每个 run 独立成段（多行）。"""
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "ctr": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    algn = {"l": PP_ALIGN.LEFT, "ctr": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = algn
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        # 中文 ea 字体，避免宋体默认发虚。
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {"typeface": "Microsoft YaHei"})
        rPr.append(ea)
    return tb


def add_rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line_color=None, alpha=None):
    sp = slide.shapes.add_shape(shape, emu(x), emu(y), emu(w), emu(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if alpha is not None:
        # 透明度：往 solidFill 里塞 <a:alpha>（python-pptx 无高层 API）。
        sf = sp.fill._xPr.find(qn("a:solidFill"))
        clr = sf.find(qn("a:srgbClr"))
        a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
        clr.append(a)
    if line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_image_cover(slide, x, y, w, h, path):
    """cover 语义：等比缩放填满框，溢出用 crop 裁掉（对齐 CSS object-fit:cover）。"""
    from PIL import Image  # python-pptx 依赖 Pillow，通常已在
    try:
        iw, ih = Image.open(path).size
    except Exception:  # noqa: BLE001
        iw = ih = 0
    pic = slide.shapes.add_picture(path, emu(x), emu(y), emu(w), emu(h))
    if iw and ih:
        target, src = w / h, iw / ih
        if src > target:
            crop = (1 - target / src) / 2
            pic.crop_left = pic.crop_right = crop
        elif src < target:
            crop = (1 - src / target) / 2
            pic.crop_top = pic.crop_bottom = crop
    return pic


def render_slide(slide, sl, pal, warnings, page):
    bg1, bg2, ink, muted, accent, card, line = [RGBColor.from_string(c) for c in pal]
    inkC, mutedC, accentC = ink, muted, accent
    layout = sl.get("layout", "bullets")

    def header(title):
        if title:
            add_text(slide, 80, 50, 1120, 64, [(title, 26, inkC, True, False)])
            add_rect(slide, 80, 122, 72, 4, accentC)

    if layout in ("title", "closing"):
        k = sl.get("kicker", "")
        if k:
            add_text(slide, 80, 210, 1120, 32, [(k, 14, accentC, True, False)], align="ctr")
        add_text(slide, 80, 260, 1120, 120, [(sl.get("title", ""), 44, inkC, True, False)], align="ctr", anchor="ctr")
        add_rect(slide, 598, 392, 84, 4, accentC)
        sub = sl.get("subtitle", "")
        if sub:
            add_text(slide, 160, 420, 960, 80, [(sub, 18, mutedC, False, False)], align="ctr")

    elif layout == "section":
        add_text(slide, 80, 300, 1120, 120, [(sl.get("title", ""), 40, inkC, True, False)], align="ctr", anchor="ctr")

    elif layout == "quote":
        add_text(slide, 100, 120, 200, 130, [("“", 96, accentC, True, False)])
        add_text(slide, 160, 250, 960, 220, [(sl.get("text", ""), 26, inkC, False, True)], align="ctr", anchor="ctr")
        by = sl.get("by", "")
        if by:
            add_text(slide, 160, 490, 960, 40, [(f"—— {by}", 15, mutedC, False, False)], align="ctr")

    elif layout == "freeform":
        boxes = sl.get("boxes") or []
        if not boxes:
            warnings.append(f"第 {page} 页 freeform 无 boxes")
        for b in boxes:
            t = b.get("type")
            x, y = b.get("x", 0), b.get("y", 0)
            w, h = max(1, b.get("w", 100)), max(1, b.get("h", 100))
            if t == "text":
                color = parse_color(b.get("color", ""), pal, pal[2])
                size = min(400, max(4, int(b.get("size", 18))))
                bold, italic = bool(b.get("bold")), bool(b.get("italic"))
                align = {"center": "ctr", "ctr": "ctr", "c": "ctr", "right": "r", "r": "r"}.get(b.get("align", ""), "l")
                anchor = {"middle": "ctr", "center": "ctr", "ctr": "ctr", "bottom": "b", "b": "b"}.get(b.get("anchor", ""), "t")
                lines = b.get("lines")
                runs = [(ln, size, color, bold, italic) for ln in lines] if isinstance(lines, list) else [(b.get("text", ""), size, color, bold, italic)]
                add_text(slide, x, y, w, h, runs, align=align, anchor=anchor)
            elif t in ("rect", "bar"):
                add_rect(slide, x, y, w, h, parse_color(b.get("color", ""), pal, pal[4]))
            elif t == "card":
                add_rect(slide, x, y, w, h, card, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=line)
            elif t == "scrim":
                add_rect(slide, x, y, w, h, parse_color(b.get("color", ""), pal, "000000"),
                         alpha=min(100, max(0, int(b.get("alpha", 50)))))
            elif t in ("image", "pic"):
                p = (b.get("image") or "").strip()
                if p:
                    try:
                        add_image_cover(slide, x, y, w, h, p)
                    except Exception as e:  # noqa: BLE001
                        warnings.append(f"第 {page} 页 freeform 图不可用: {e}")
                else:
                    warnings.append(f"第 {page} 页 freeform 图框缺 image 路径")
            else:
                warnings.append(f"第 {page} 页 freeform 未知盒子 \"{t}\"")

    elif layout in ("image-full", "image-text"):
        img = (sl.get("image") or "").strip()
        ok = False
        if img:
            try:
                if layout == "image-full":
                    add_image_cover(slide, 0, 0, CANVAS_W, CANVAS_H, img)
                    add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, RGBColor.from_string("000000"), alpha=50)
                    add_text(slide, 80, 268, 1120, 110, [(sl.get("title", ""), 40, RGBColor.from_string("FFFFFF"), True, False)], align="ctr", anchor="ctr")
                    sub = sl.get("subtitle", "")
                    if sub:
                        add_text(slide, 160, 420, 960, 70, [(sub, 17, RGBColor.from_string("E8E8E8"), False, False)], align="ctr")
                else:
                    header(sl.get("title", ""))
                    right = str(sl.get("side", "")).lower() == "right"
                    ix, tx = (656, 80) if right else (80, 656)
                    add_image_cover(slide, ix, 168, 544, 470, img)
                    runs = [(p if isinstance(p, str) else p.get("text", ""), 16, inkC, False, False) for p in (sl.get("points") or [])]
                    if runs:
                        add_text(slide, tx, 180, 544, 446, runs)
                ok = True
            except Exception as e:  # noqa: BLE001
                warnings.append(f"第 {page} 页配图不可用: {e}")
        if not ok:
            warnings.append(f"第 {page} 页 {layout} 缺可用配图，已降级标题")
            add_text(slide, 80, 260, 1120, 120, [(sl.get("title", ""), 40, inkC, True, False)], align="ctr", anchor="ctr")

    else:  # bullets / two-col / compare / stats / timeline / 未知 → 通用要点渲染
        if layout not in ("bullets", "two-col", "compare", "stats", "timeline") and layout:
            warnings.append(f"第 {page} 页未知版式 \"{layout}\"，按 bullets 渲染")
        header(sl.get("title", ""))
        runs = []
        for p in (sl.get("points") or []):
            txt = p if isinstance(p, str) else p.get("text", "")
            runs.append(("• " + txt, 17, inkC, False, False))
            if isinstance(p, dict):
                for s in (p.get("sub") or []):
                    runs.append(("    – " + s, 14, mutedC, False, False))
        if runs:
            add_text(slide, 80, 176, 1120, 470, runs)


def main():
    if len(sys.argv) < 3:
        sys.stderr.write("用法: pptx_bridge.py <spec.json> <out.pptx>\n")
        sys.exit(2)
    spec_path, out_path = sys.argv[1], sys.argv[2]
    with open(spec_path, "r", encoding="utf-8-sig") as f:
        spec = json.load(f)
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        sys.stderr.write("spec 缺 slides 数组或为空\n")
        sys.exit(4)

    theme = spec.get("theme", "")
    pal = palette(theme)
    warnings = []

    prs = Presentation()
    prs.slide_width = emu(CANVAS_W)
    prs.slide_height = emu(CANVAS_H)
    blank = prs.slide_layouts[6]  # 完全空白版式

    img_count = 0
    for i, sl in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, pal)
        before = len(slide.shapes)
        render_slide(slide, sl, pal, warnings, i + 1)
        img_count += sum(1 for sh in slide.shapes if sh.shape_type == 13)  # 13 = PICTURE
        notes = (sl.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(out_path)
    print(json.dumps({
        "ok": True, "out": out_path, "slides": len(slides),
        "theme": theme, "images": img_count, "editable": True,
        "engine": "python", "warnings": warnings,
    }, ensure_ascii=False))


if __name__ == "__main__":
    # 想复用 build/engine.py 的 THEMES/版式：解开下面两行，把 spec 映射到 engine 方法即可。
    #   sys.path.insert(0, "/mnt/d/polaris/教师助手/build"); import engine
    main()
