# -*- coding: utf-8 -*-
"""版式构件库：统一页眉/页脚/正文区 + 真实字宽测量的自动折行与溢出断言。

设计前提（goal 第四、六节）：
  · 所有文本框都留内边距，文字不贴框；
  · 文字高度用 PIL 按真实字体量出来，放不下就报错，绝不靠缩小字号硬塞；
  · 图片一律等比，place 时宽高同一比例因子；
  · 正文不得越过页脚安全区 FOOTER_Y。
"""
import os

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from theme import (INK, INK_SOFT, MUTED, CYAN, CORAL, AMBER, BG, BG2, RULE, WHITE,
                   FONT, SLIDE_W_IN, SLIDE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
                   TITLE_Y, TITLE_H, RULE_Y, BODY_TOP, BODY_BOTTOM, FOOTER_Y,
                   PT_TITLE, PT_H2, PT_BODY, PT_SMALL, PT_FOOTER, PT_MIN_BODY,
                   PAD_L, PAD_R, PAD_T, PAD_B)

# ───────────────────────── 字宽测量 ─────────────────────────
# 避头尾：这些标点不许出现在行首（中文排版惯例），否则会看到孤零零的一个「。」占一整行
_NO_LINE_START = set("。，、；：？！）】》」』”’%》…·")

_FONT_FILES = [
    r"C:\Windows\Fonts\msyh.ttc",
    r"C:\Windows\Fonts\msyhbd.ttc",
]
_CACHE = {}


def set_run_font(run, name):
    """同时设置 latin / ea / cs 三个字面。

    python-pptx 的 font.name 只写 <a:latin>，中文会落到主题的东亚字体（等线），
    字宽与我们用来测量的 msyh 不一致 → 排版按 msyh 算、PowerPoint 按等线排，
    结果就是「静态检查通过、实机却溢出卡片」。必须把 ea 也钉死。
    """
    from pptx.oxml.ns import qn
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            (latin.addnext(el) if latin is not None else rPr.append(el))
        el.set("typeface", name)
        latin = el


def _font(pt, bold=False):
    """按 4× 精度加载字体用于测量；返回 (font, px_per_pt)。"""
    key = (round(pt, 1), bold)
    if key not in _CACHE:
        path = _FONT_FILES[1] if bold else _FONT_FILES[0]
        _CACHE[key] = ImageFont.truetype(path, int(round(pt * 4)))
    return _CACHE[key], 4.0


def text_w_pt(s, pt, bold=False):
    """字符串在给定字号下的排版宽度（磅）。"""
    f, k = _font(pt, bold)
    return f.getlength(s) / k


def wrap(text, pt, max_w_in, bold=False):
    """按真实字宽折行。中文逐字断，英文/数字整词不拆。"""
    limit = max_w_in * 72.0
    out = []
    for para in str(text).split("\n"):
        if not para:
            out.append("")
            continue
        # 切成「不可拆的最小单元」：CJK/标点 单字，ASCII 连续串成词
        toks, buf = [], ""
        for ch in para:
            if ord(ch) < 128 and not ch.isspace():
                buf += ch
            else:
                if buf:
                    toks.append(buf); buf = ""
                toks.append(ch)
        if buf:
            toks.append(buf)
        cur = ""
        for t in toks:
            cand = cur + t
            # +0.01pt 容差：inner_w 由浮点相减得到（0.9−0.2−0.2=0.49999…），
            # 不留容差会让「刚好占满一行」的短标签在 build 与 check 两侧算出不同行数。
            if cur and text_w_pt(cand, pt, bold) > limit + 0.01:
                # 注意：这里刻意**不做**避头尾（标点悬挂）。PowerPoint 在中英混排处
                # 会把「2。」拆成「2」+「。」，我们若按悬挂算就会少算一行 → 实机溢出。
                # 宁可高估一行，也不能低估。
                out.append(cur.rstrip())
                cur = "" if t.isspace() else t
            else:
                cur = cand
        out.append(cur.rstrip())
    return out


def block_h_in(lines, pt, spacing=1.35, space_after_pt=6.0):
    """一段文本（已折行）的实际占高（英寸）。PowerPoint 单倍行距 ≈ 1.2×字号。"""
    if not lines:
        return 0.0
    n = len(lines)
    return (n * pt * 1.2 * spacing + max(0, n - 1) * 0 + space_after_pt) / 72.0


# ───────────────────────── Deck ─────────────────────────
class Deck:
    def __init__(self, footer_left="", pool=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W_IN)
        self.prs.slide_height = Inches(SLIDE_H_IN)
        self.footer_left = footer_left
        self.pool = pool
        self.pages = []          # [(slide, 页面用途)]

    # ── 基础 ──
    def blank(self, purpose="", bg=BG, footer=True):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb = bg
        r.line.fill.background(); r.shadow.inherit = False
        r.name = "bg_page"
        self.pages.append((s, purpose))
        if footer:
            self.footer(s, len(self.pages))
        return s

    def footer(self, slide, no, right=None):
        a = self.text(slide, MARGIN_L, FOOTER_Y + 0.10, CONTENT_W * 0.7, 0.3,
                      self.footer_left, pt=PT_FOOTER, color=MUTED, spacing=1.0)
        b = self.text(slide, SLIDE_W_IN - MARGIN_R - 1.6, FOOTER_Y + 0.10, 1.6, 0.3,
                      right if right is not None else f"{no:02d}", pt=PT_FOOTER, color=MUTED,
                      align="r", spacing=1.0)
        a.name, b.name = "footer_left", "footer_right"
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(FOOTER_Y),
                                    Inches(CONTENT_W), Emu(9525))
        ln.fill.solid(); ln.fill.fore_color.rgb = RULE
        ln.line.fill.background(); ln.shadow.inherit = False

    # ── 文本 ──
    def text(self, slide, x, y, w, h, text, pt=PT_BODY, color=INK, bold=False,
             align="l", anchor="t", spacing=1.35, font=FONT, space_after=6,
             wrap_text=True, check=True, grow=True, role="body", name=None):
        """放一段文字。按真实字宽折行；框高不够时默认把框**撑高**（grow）而不是让文字溢出，
        check=True 且 grow=False 时改为直接报错。相邻元素的碰撞交给 check.py 的重叠检测兜底。"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        # 形状名带 role 前缀：check.py 据此分层判字号（正文≥20pt，标签/图注≥15pt）
        tb.name = name or f"{role}_{len(slide.shapes)}"
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left, tf.margin_right = PAD_L, PAD_R
        tf.margin_top, tf.margin_bottom = PAD_T, PAD_B
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "ctr": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[anchor]
        inner_w = w - (PAD_L.inches + PAD_R.inches)
        lines = wrap(text, pt, inner_w, bold) if wrap_text else str(text).split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"l": PP_ALIGN.LEFT, "ctr": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
            p.line_spacing = spacing
            p.space_after = Pt(space_after)
            run = p.add_run(); run.text = ln
            f = run.font
            f.size = Pt(pt); f.bold = bold; f.color.rgb = color
            set_run_font(run, font)
        need = block_h_in(lines, pt, spacing, space_after) + PAD_T.inches + PAD_B.inches
        if need > h + 0.02:
            if grow:
                tb.height = Inches(need)
            elif check:
                raise LayoutError(
                    f"文本超框：需 {need:.2f}in > 框高 {h:.2f}in\n  «{str(text)[:60]}…»\n"
                    f"  （{len(lines)} 行 @{pt}pt）请加高框、减字数或分页——不许缩字号。")
        return tb

    def measure_text_h(self, text, pt, w, spacing=1.35, space_after=6, bold=False):
        inner_w = w - (PAD_L.inches + PAD_R.inches)
        return block_h_in(wrap(text, pt, inner_w, bold), pt, spacing, space_after) \
            + PAD_T.inches + PAD_B.inches

    # ── 自动高度的两个便捷件（内容脚本几乎只用这两个）──
    def para(self, slide, x, y, w, text, pt=PT_BODY, color=INK, bold=False, align="l",
             spacing=1.30, role="body"):
        """正文段落：高度按真实字宽算好再建框，返回占高，调用方直接累加 y。"""
        h = self.measure_text_h(text, pt, w, spacing=spacing, space_after=0, bold=bold)
        self.text(slide, x, y, w, h, text, pt=pt, color=color, bold=bold, align=align,
                  spacing=spacing, space_after=0, role=role)
        return h

    def label(self, slide, x, y, w, text, pt=15, color=AMBER, bold=True, align="l",
              role="tag"):
        """小标签/图注：同样自量高度。role=tag/cap 的字号下限比正文低。"""
        return self.para(slide, x, y, w, text, pt=pt, color=color, bold=bold, align=align,
                         spacing=1.0, role=role)

    # ── 页眉 ──
    def header(self, slide, title, kicker=None, tag=None, tag_color=CYAN):
        if kicker:
            self.text(slide, MARGIN_L, TITLE_Y - 0.28, CONTENT_W * 0.6, 0.30, kicker,
                      pt=15, color=CYAN, bold=True, spacing=1.0, space_after=0, role="tag")
            ty = TITLE_Y + 0.08
        else:
            ty = TITLE_Y
        self.text(slide, MARGIN_L, ty, CONTENT_W - (2.4 if tag else 0), TITLE_H, title,
                  pt=PT_TITLE, color=INK, bold=True, spacing=1.1, space_after=0, anchor="t",
                  role="h1")
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(RULE_Y),
                                    Inches(0.86), Inches(0.045))
        ln.fill.solid(); ln.fill.fore_color.rgb = AMBER
        ln.line.fill.background(); ln.shadow.inherit = False
        ln2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN_L + 0.86), Inches(RULE_Y + 0.017),
                                     Inches(CONTENT_W - 0.86), Emu(9525))
        ln2.fill.solid(); ln2.fill.fore_color.rgb = RULE
        ln2.line.fill.background(); ln2.shadow.inherit = False
        if tag:
            tw = max(1.5, text_w_pt(tag, 15, True) / 72.0 + 0.55)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(SLIDE_W_IN - MARGIN_R - tw), Inches(TITLE_Y + 0.05),
                                         Inches(tw), Inches(0.44))
            box.name = "tag_chip"
            box.fill.solid(); box.fill.fore_color.rgb = tag_color
            box.line.fill.background(); box.shadow.inherit = False
            box.adjustments[0] = 0.28
            tf = box.text_frame; tf.word_wrap = False
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = tag
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
            set_run_font(r, FONT)

    # ── 图形件 ──
    def card(self, slide, x, y, w, h, fill=BG2, line=RULE, bar=None, radius=0.06):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                     Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = fill
        if line is not None:
            box.line.color.rgb = line; box.line.width = Pt(1.0)
        else:
            box.line.fill.background()
        box.shadow.inherit = False
        try:
            box.adjustments[0] = radius
        except Exception:
            pass
        if bar is not None:
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.075), Inches(h))
            b.fill.solid(); b.fill.fore_color.rgb = bar
            b.line.fill.background(); b.shadow.inherit = False
        return box

    def rule(self, slide, x, y, w, color=RULE, thick=0.013):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(thick))
        r.fill.solid(); r.fill.fore_color.rgb = color
        r.line.fill.background(); r.shadow.inherit = False
        return r

    def vrule(self, slide, x, y, h, color=RULE, thick=0.013):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(thick), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = color
        r.line.fill.background(); r.shadow.inherit = False
        return r

    def bullet_line(self, slide, x, y, w, text, pt=PT_BODY, color=INK, dot=CYAN, bold=False):
        """带圆点的一行/多行要点，返回占高。"""
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + pt * 1.2 / 72 * 0.32),
                                   Inches(0.13), Inches(0.13))
        d.fill.solid(); d.fill.fore_color.rgb = dot
        d.line.fill.background(); d.shadow.inherit = False
        h = self.measure_text_h(text, pt, w - 0.28)
        self.text(slide, x + 0.26, y - 0.06, w - 0.26, h, text, pt=pt, color=color, bold=bold)
        return h

    # ── 图片（等比） ──
    def picture(self, slide, path, x, y, box_w, box_h, mode="contain", align="ctr"):
        """等比放置图片。contain=完整显示（居中留白）；cover=等比裁切填满。"""
        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
        ar = iw / ih
        if mode == "contain":
            w, h = (box_w, box_w / ar) if box_w / ar <= box_h else (box_h * ar, box_h)
            ox = x + (box_w - w) / 2 if align == "ctr" else x
            oy = y + (box_h - h) / 2
            return slide.shapes.add_picture(path, Inches(ox), Inches(oy), Inches(w), Inches(h))
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
        box_ar = box_w / box_h
        if ar > box_ar:                       # 原图更宽 → 左右裁
            c = (1 - box_ar / ar) / 2
            pic.crop_left = pic.crop_right = c
        elif ar < box_ar:                     # 原图更高 → 上下裁
            c = (1 - ar / box_ar) / 2
            pic.crop_top = pic.crop_bottom = c
        return pic

    def notes(self, slide, text):
        slide.notes_slide.notes_text_frame.text = text

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        try:
            self.prs.save(path)
        except PermissionError:
            alt = path[:-5] + "-新.pptx"
            self.prs.save(alt)
            print("  (原文件被占用) ->", alt)
            return alt
        return path


class LayoutError(RuntimeError):
    pass
